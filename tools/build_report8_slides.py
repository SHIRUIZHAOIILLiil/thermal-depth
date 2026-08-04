"""Build the report-8 deck from docs/AIRE_RESULTS_20260802.md.

Every number here is transcribed from that document, which is the single source
of truth for the Aire run; nothing is recomputed and nothing is rounded further.
The deck exists as a script rather than as a hand-edited pptx because report 6's
build script was not kept, and the numbers had to be re-entered by hand for
report 7.

Storyline (the supervisor-facing spine, in this order):
    six-route baseline -> caption's two opposite effects -> d2's far-field
    advantage under a controlled attribution -> the external reference point
    (original AnyThermal) -> limits.
Rain and gradient matching are supplements, not spine: the first is a condition
boundary on the caption finding, the second is a clean negative result.

Each of the six routes gets a data-flow slide, drawn with the shape vocabulary
report 7 established for its route-B figure (same trapezoids, cubes, badges and
palette) so the two decks read as one series.  The rule the legend states:
SHAPE says what a block is, COLOUR says whether it takes gradients.  Module
names and parameter counts come from docs/ROUTE_ARCHITECTURES_20EPOCH.md, which
tools/dump_route_architecture.py generates from the models themselves.

One slide carries a placeholder on purpose -- the qualitative comparison is a
cluster job that had not returned when the deck was built (docs §11).  It is
marked in the deck itself, not silently omitted, so an unfilled box is visible
rather than forgotten.

    python tools/build_report8_slides.py --output "E:/.../8-thermal_depth_report_8.pptx"
"""

from __future__ import annotations

import argparse
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

INK = RGBColor(0x1F, 0x2A, 0x44)      # deep charcoal blue -- titles, dark bars
ICE = RGBColor(0xCA, 0xDC, 0xFC)      # ice blue -- kickers on dark, accents
GREY = RGBColor(0x5A, 0x66, 0x7A)     # muted -- kickers on light, footnotes
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOOD = RGBColor(0x1E, 0x6B, 0x52)     # an improvement
BAD = RGBColor(0xA6, 0x39, 0x2E)      # a regression
RULE = RGBColor(0xD5, 0xDB, 0xE4)
BAND = RGBColor(0xF2, 0xF5, 0xF9)     # table header / zebra fill

FONT = "微软雅黑"
DATE = "2026-08-03"

W, H = 13.333, 7.5
MARGIN, BODY_W = 0.60, 12.13
BAR_TOP, BAR_H = 6.62, 0.66

# ── data-flow diagram palette ───────────────────────────────────────────────
# Report 7's route-B figure established these; reused verbatim so the two decks
# read as one series.  The rule the legend states: SHAPE says what a block is,
# COLOUR says whether it takes gradients.
VAE_FILL, VAE_LINE = RGBColor(0xC7, 0xDC, 0xF0), RGBColor(0x6E, 0x9C, 0xC4)
TXT_FILL, TXT_LINE = RGBColor(0xD6, 0xE8, 0xCE), RGBColor(0x8F, 0xB0, 0x7F)
ANY_FILL, ANY_LINE = RGBColor(0xE2, 0xDC, 0xF2), RGBColor(0x91, 0x87, 0xC4)
HOT_FILL, HOT_LINE = RGBColor(0xE8, 0x8B, 0x2F), RGBColor(0xC0, 0x6E, 0x18)
COLD_FILL, COLD_LINE = RGBColor(0xD8, 0xDE, 0xE8), RGBColor(0x8A, 0x96, 0xAA)
CUBE_FILL, CUBE_LINE = RGBColor(0xEC, 0xE6, 0xD6), RGBColor(0xB6, 0xAC, 0x96)
WIRE = RGBColor(0x7C, 0x89, 0x9E)
BADGE_COLD, BADGE_HOT = RGBColor(0x8A, 0x96, 0xAA), RGBColor(0xC0, 0x6E, 0x18)
PROMPT_FILL, PROMPT_LINE = RGBColor(0xFD, 0xF3, 0xDD), RGBColor(0xD8, 0xB6, 0x6A)

FIGDIR = Path(__file__).resolve().parents[1] / "docs" / "figures" / "routeflow"

ROW_Y, TEXT_Y = 2.86, 4.86      # centre lines of the main path and the text branch
GAP = 0.30                      # horizontal gap between consecutive nodes


# ── primitives ──────────────────────────────────────────────────────────────

def textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    return frame


def write(frame, lines, *, size, colour=INK, bold=False, space_after=4, line_spacing=None):
    """`lines` is a list of str, or of (text, {overrides}) pairs."""
    for index, line in enumerate(lines):
        text, over = line if isinstance(line, tuple) else (line, {})
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.space_after = Pt(over.get("space_after", space_after))
        if line_spacing:
            para.line_spacing = line_spacing
        run = para.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(over.get("size", size))
        run.font.bold = over.get("bold", bold)
        run.font.color.rgb = over.get("colour", colour)
    return frame


def header(slide, kicker, title):
    write(textbox(slide, MARGIN, 0.38, 12.00, 0.30), [kicker], size=11, colour=GREY, bold=True)
    write(textbox(slide, MARGIN, 0.68, 12.10, 0.72), [title], size=26, colour=INK, bold=True)


def conclusion(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(BAR_TOP),
                                 Inches(BODY_W), Inches(BAR_H))
    bar.fill.solid()
    bar.fill.fore_color.rgb = INK
    bar.line.fill.background()
    bar.shadow.inherit = False
    frame = bar.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.18)
    frame.margin_right = Inches(0.18)
    para = frame.paragraphs[0]
    prefix = para.add_run()
    prefix.text = "结论   "
    prefix.font.name, prefix.font.size, prefix.font.bold = FONT, Pt(12), True
    prefix.font.color.rgb = ICE
    body = para.add_run()
    body.text = text
    body.font.name, body.font.size, body.font.bold = FONT, Pt(12), False
    body.font.color.rgb = WHITE


def footnote(slide, text, top=6.24):
    write(textbox(slide, MARGIN, top, BODY_W, 0.32), [text], size=10.5, colour=GREY)


def cell_text(cell, text, *, size, bold=False, colour=INK, align=PP_ALIGN.LEFT):
    frame = cell.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = frame.margin_right = Inches(0.06)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name, run.font.size, run.font.bold = FONT, Pt(size), bold
    run.font.color.rgb = colour


def table(slide, rows, left, top, width, height, col_widths, *, size=11.5,
          align=None, highlight=()):
    """`rows[0]` is the header. A cell may be a str or (str, {overrides})."""
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tbl = shape.table
    tbl.first_row = True
    tbl.horz_banding = False
    for index, fraction in enumerate(col_widths):
        tbl.columns[index].width = Inches(width * fraction / sum(col_widths))
    for r, row in enumerate(rows):
        for c, raw in enumerate(row):
            text, over = raw if isinstance(raw, tuple) else (raw, {})
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BAND if (r == 0 or r in highlight) else WHITE
            cell_text(cell, text,
                      size=over.get("size", size),
                      bold=over.get("bold", r == 0 or r in highlight),
                      colour=over.get("colour", INK if r == 0 else INK),
                      align=over.get("align",
                                     (align[c] if align else PP_ALIGN.LEFT) if c else PP_ALIGN.LEFT))
    return tbl


def placeholder(slide, left, top, width, height, lines):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = BAND
    box.line.color.rgb = RULE
    box.line.width = Pt(1.25)
    box.shadow.inherit = False
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.24)
    frame.margin_top = Inches(0.20)
    write(frame, lines, size=12, colour=GREY)


# ── data-flow diagram ───────────────────────────────────────────────────────

def _shape(slide, kind, cx, cy, vw, vh, *, rot=0, fill=None, line=None, lw=1.25):
    """Place `kind` by its VISUAL centre and visual width/height.

    PowerPoint rotates about the centre and leaves left/top/width/height
    describing the UNROTATED box, so a 90-degree trapezoid needs its w/h
    swapped before placing.  Doing that here keeps every call site in visual
    coordinates, which is the only way the layout arithmetic stays readable.
    """
    bw, bh = (vh, vw) if rot in (90, 270) else (vw, vh)
    shape = slide.shapes.add_shape(kind, Inches(cx - bw / 2), Inches(cy - bh / 2),
                                   Inches(bw), Inches(bh))
    shape.rotation = rot
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    return shape


def _centred(slide, cx, cy, width, lines, *, size, colour=INK):
    frame = textbox(slide, cx - width / 2, cy, width, 0.5)
    write(frame, lines, size=size, colour=colour, space_after=0)
    for para in frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
    return frame


def _badge(slide, cx, top, trained):
    box = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx, top + 0.105, 0.46, 0.21,
                 fill=BADGE_HOT if trained else BADGE_COLD, line=None)
    frame = box.text_frame
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = "训练" if trained else "冻结"
    run.font.name, run.font.size, run.font.bold = FONT, Pt(8), True
    run.font.color.rgb = WHITE


def _wire(slide, x1, y1, x2, y2):
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = WIRE
    connector.line.width = Pt(1.5)
    return connector


def _node(slide, node, cx, cy):
    """Draw one block; returns its visual width."""
    kind = node["kind"]
    vw, vh = node["w"], node.get("h", 1.20)
    trained = node.get("trained", False)

    if kind == "image":
        image = FIGDIR / node["file"]
        if not image.is_file():
            # A route whose prediction panel has not been exported yet: show the gap
            # rather than borrowing another route's output, which would silently put
            # the wrong scene next to the right input.
            _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, vw, vw * 0.4,
                   fill=BAND, line=RULE, lw=1.0)
            _centred(slide, cx, cy - 0.10, vw, [f"待填\n{node['file']}"], size=8.5, colour=GREY)
            _centred(slide, cx, cy + vw * 0.2 + 0.06, vw + 0.9,
                     [node["label"], (node["sub"], {"size": 7.5})], size=9.5, colour=GREY)
            return vw
        picture = slide.shapes.add_picture(str(image), 0, 0, width=Inches(vw))
        # picture.height is already EMU; Inches() would read it as a count of
        # inches and throw the picture hundreds of feet off the slide
        drawn_h = Emu(picture.height).inches
        picture.left = Inches(cx - vw / 2)
        picture.top = Inches(cy - drawn_h / 2)
        _centred(slide, cx, cy + drawn_h / 2 + 0.06, vw + 0.9,
                 [node["label"], (node["sub"], {"size": 7.5})], size=9.5, colour=GREY)
        return vw

    if kind == "cube":
        _shape(slide, MSO_SHAPE.CUBE, cx, cy, vw, vh, fill=CUBE_FILL, line=CUBE_LINE, lw=1.0)
        _centred(slide, cx, cy + vh / 2 + 0.02, 1.60,
                 [node["label"], (node["sub"], {"size": 7.5})], size=9.5)
        return vw

    if kind in ("enc", "dec"):
        fill, line = node.get("fill", VAE_FILL), node.get("line", VAE_LINE)
        _shape(slide, MSO_SHAPE.TRAPEZOID, cx, cy, vw, vh,
               rot=90 if kind == "enc" else 270, fill=fill, line=line)
        _centred(slide, cx, cy - 0.26, max(1.75, vw + 0.7),
                 [node["label"], (node["sub"], {"size": 8})], size=10)
        _badge(slide, cx - min(0.33, vw / 2 - 0.05), cy - vh / 2 + 0.02, trained)
        return vw

    if kind == "unet":                      # hourglass = two facing trapezoids
        fill, line = (HOT_FILL, HOT_LINE) if trained else (COLD_FILL, COLD_LINE)
        half = vw / 2
        _shape(slide, MSO_SHAPE.TRAPEZOID, cx - half / 2, cy, half, vh, rot=90,
               fill=fill, line=line)
        _shape(slide, MSO_SHAPE.TRAPEZOID, cx + half / 2, cy, half, vh, rot=270,
               fill=fill, line=line)
        _centred(slide, cx, cy - 0.26, vw + 0.9,
                 [node["label"], (node["sub"], {"size": 8})], size=10,
                 colour=WHITE if trained else INK)
        _badge(slide, cx - 0.33, cy - vh / 2 + 0.02, trained)
        return vw

    if kind == "block":                     # adapter: a residual CNN, not an encoder
        fill, line = (HOT_FILL, HOT_LINE) if trained else (COLD_FILL, COLD_LINE)
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, vw, vh, fill=fill, line=line)
        _centred(slide, cx, cy - 0.26, max(1.75, vw + 0.6),
                 [node["label"], (node["sub"], {"size": 8})], size=10,
                 colour=WHITE if trained else INK)
        _badge(slide, cx - min(0.33, vw / 2 - 0.05), cy - vh / 2 + 0.02, trained)
        return vw

    if kind == "concat":
        _shape(slide, MSO_SHAPE.OVAL, cx, cy, 0.44, 0.44, fill=WHITE, line=WIRE)
        _centred(slide, cx, cy - 0.16, 0.6, ["⊕"], size=14)
        _centred(slide, cx, cy - 0.78, 1.6, ["concat → 8 通道"], size=8.5, colour=GREY)
        return 0.44

    raise ValueError(f"unknown node kind {kind!r}")


def flow(slide, nodes, *, loss="掩码 SSI-L1 vs LiDAR 视差", prompt='c = ""（空 prompt）',
         prompts=None):
    """Lay the main path out left to right, then hang the text branch under it.

    `prompts` stacks several alternative prompt boxes into the one text branch.
    The caption figure needs that: its whole point is that the image path is
    fixed and only the string entering CLIP changes.
    """
    total = sum(n["w"] for n in nodes) + GAP * (len(nodes) - 1)
    x = MARGIN + 0.45 + (BODY_W - 0.9 - total) / 2
    centres = []
    for node in nodes:
        cx = x + node["w"] / 2
        _node(slide, node, cx, ROW_Y)
        centres.append((cx, node["w"]))
        x += node["w"] + GAP
    for (cx, w), (nx, nw) in zip(centres, centres[1:]):
        _wire(slide, cx + w / 2 + 0.02, ROW_Y, nx - nw / 2 - 0.02, ROW_Y)

    concat_i = next(i for i, n in enumerate(nodes) if n["kind"] == "concat")
    unet_i = next(i for i, n in enumerate(nodes) if n["kind"] == "unet")
    concat_x = centres[concat_i][0]
    unet_x = centres[unet_i][0]

    # text branch: prompt -> CLIP -> cross-attention into the U-Net from below
    px = MARGIN + 0.75
    boxes = prompts or [prompt]
    span = 0.50
    # the stack sits a little above the branch line so its last box clears the
    # body text underneath; the wires still converge on the encoder's own centre
    top = (TEXT_Y - 0.12) - span * (len(boxes) - 1) / 2
    for index, text in enumerate(boxes):
        y = top + index * span
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, px + 0.72, y, 1.62, 0.40,
               fill=PROMPT_FILL, line=PROMPT_LINE)
        _centred(slide, px + 0.72, y - 0.09, 1.72, [text], size=8.5)
        if len(boxes) > 1:
            _wire(slide, px + 1.53, y, px + 1.86, TEXT_Y)
    clip_x = px + 2.42
    _shape(slide, MSO_SHAPE.TRAPEZOID, clip_x, TEXT_Y, 0.84, 1.15, rot=90,
           fill=TXT_FILL, line=TXT_LINE)
    _centred(slide, clip_x, TEXT_Y - 0.26, 1.75,
             ["Text Encoder", ("340.39 M", {"size": 8})], size=10)
    _badge(slide, clip_x - 0.33, TEXT_Y - 0.575 + 0.02, False)
    _wire(slide, px + 1.86 if len(boxes) > 1 else px + 1.53, TEXT_Y,
          clip_x - 0.44, TEXT_Y)
    _wire(slide, clip_x + 0.44, TEXT_Y, unet_x, TEXT_Y)
    _wire(slide, unet_x, TEXT_Y, unet_x, ROW_Y + 0.55)
    _centred(slide, unet_x + 1.05, TEXT_Y - 0.44, 1.7, ["cross-attention"],
             size=8.5, colour=GREY)

    # Noise latent joins at the concat. It sits well above the text row: its
    # caption would otherwise land exactly on the wire running to the U-Net.
    noise_y = ROW_Y + 1.00
    _shape(slide, MSO_SHAPE.CUBE, concat_x, noise_y, 0.58, 0.68,
           fill=CUBE_FILL, line=CUBE_LINE, lw=1.0)
    _centred(slide, concat_x, noise_y + 0.36, 1.7,
             ["z_T ~ N(0, I)", ("4 × 32 × 80", {"size": 7.5})], size=9.5)
    _wire(slide, concat_x, noise_y - 0.34, concat_x, ROW_Y + 0.24)

    # loss box hangs off the decoded output
    out_x = centres[-1][0]
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, out_x, TEXT_Y - 0.10, 2.10, 0.46,
           fill=INK, line=None)
    _centred(slide, out_x, TEXT_Y - 0.20, 2.10, [loss], size=9, colour=WHITE)
    _wire(slide, out_x, ROW_Y + 0.62, out_x, TEXT_Y - 0.33)


def picture_or_placeholder(slide, image: Path | None, left, top, width, height, lines):
    """Fit the picture inside the box and centre it, rather than forcing width.

    The comparison strip is about 2.1:1; setting only the width would make it
    5.7in tall in a 4.5in slot and push it through the conclusion bar.
    """
    if not (image and image.is_file()):
        placeholder(slide, left, top, width, height, lines)
        return
    with Image.open(image) as probe:
        aspect = probe.size[0] / probe.size[1]
    draw_w = min(width, height * aspect)
    draw_h = draw_w / aspect
    slide.shapes.add_picture(str(image),
                             Inches(left + (width - draw_w) / 2),
                             Inches(top + (height - draw_h) / 2),
                             width=Inches(draw_w))


# ── slides ──────────────────────────────────────────────────────────────────

def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    back = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    back.fill.solid()
    back.fill.fore_color.rgb = INK
    back.line.fill.background()
    back.shadow.inherit = False
    write(textbox(slide, 1.00, 2.25, 11.30, 0.40),
          ["THERMAL DEPTH · PROGRESS REPORT 8"], size=13, colour=ICE, bold=True)
    write(textbox(slide, 1.00, 2.75, 11.30, 1.20),
          ["caption 的两个相反效应，与远端误差的归因"], size=40, colour=WHITE, bold=True)
    write(textbox(slide, 1.00, 4.05, 11.30, 1.50), [
        "六条路线 × 20 epoch，全部在利兹 Aire 超算完成并在独立 test 集出数",
        "MS2 官方 BMSD 协议 · test = 16-08-46，2,543 帧，训练与调参从未接触",
    ], size=15, colour=ICE, space_after=6)
    write(textbox(slide, 1.00, 6.30, 11.30, 0.40), [DATE], size=13, colour=GREY)


def overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "WHAT THIS ROUND ESTABLISHED", "本轮拿到的四件事")
    rows = [
        ["", "发现", "证据强度"],
        ["1",
         "caption 的作用是两个方向相反的效应：训练时用它有益，推理时喂它有害。"
         "全项目最好的热像成绩来自「用 caption 训练、推理给空 prompt」＝ 0.0869。",
         "注入效应：同权重同帧、仅 prompt 不同，无懈可击；权重效应：单 seed"],
        ["2",
         "天空被判成近处，是解冻 U-Net 的代价：只有 U-Net 冻结的 c1 把上带放在 16.1 m"
         "（真值 14.8 m），解冻的 b / c2 / d2 全掉到 10–12 m。"
         "而 c1 的整体精度是六线最差 —— 两者在当前设计里对立。",
         "c1 五帧全胜，机制与外部模型一致；但只有 5 帧，且 d1 只弱支持"],
        ["3",
         "远端误差与 condition 来源有关：c2 与 d2 的远端 0.1817 → 0.1524（−16%）。"
         "但两者除 condition 来源外，Adapter 的训练方式也不同，归因只能到「二者之一」。",
         "donor-swap 证伪通过（误差翻 3.3 倍）；但对照不干净"],
        ["4",
         "远端退化不是热像深度的通病：原版 AnyThermal 用同一份稀疏 LiDAR 监督，"
         "远端几乎不退化（×1.13 vs 我方 ×1.90）。差别在架构，不在监督。",
         "外部模型，同一考卷，同一分层定义"],
    ]
    table(slide, rows, MARGIN, 1.66, BODY_W, 4.20, [0.4, 7.0, 3.8], size=11)
    write(textbox(slide, MARGIN, 6.02, BODY_W, 0.55), [
        "1 和 2 是本轮的新东西：前者解释了以往为什么总测不出 caption 的效果"
        "（两个相反效应叠在一起互相抵消），后者把「天空问题」从现象变成了可归因的设计选择。",
    ], size=11.5, colour=GREY)
    conclusion(slide, "六条线全部完成 20 epoch 并在独立 test 集出数；四件事各有可检验的证据，强度不同，逐条标注。")


def protocol(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "PROTOCOL & GOVERNANCE", "口径：所有数字共用的那一套")
    rows = [
        ["项", "值"],
        ["val（选 epoch 用）", "11-23-45，5,810 帧"],
        ["test（对外数字）", "16-08-46，2,543 帧，官方 val split、白天，训练与调参从未接触"],
        ["协议", "official BMSD ssi_disparity，min_depth 1e-3，max_depth 80"],
        ["对齐", "六线一律 ssi_disparity；原版 AnyThermal 用 ssi（官方表）"],
        ["训练 manifest", "ms2_train_day2seq_clip75_20260728.jsonl，19,949 帧"],
        ["环境", "torch 2.3.1+cu121 / L40S 48GB"],
    ]
    table(slide, rows, MARGIN, 1.66, BODY_W, 2.55, [2.6, 9.5], size=12)

    write(textbox(slide, MARGIN, 4.42, 5.90, 1.70), [
        ("⚠️ 两族模型的对齐方式不可混用", {"bold": True, "colour": BAD}),
        "把六线的预测按 ssi 对齐会得到 0.37（正确值 0.0884）；",
        "把 AnyThermal 按 ssi_disparity 对齐会得到 0.2724（正确值 0.0821）。",
        "所以不同对齐的模型必须分开跑、并排读。",
    ], size=12, colour=INK, space_after=3)

    write(textbox(slide, 6.85, 4.42, 5.88, 1.70), [
        ("✅ 跨机器标定已通过", {"bold": True, "colour": GOOD}),
        "集群 b(e05) 全量 val 0.07749105 vs 本地 0.0775，",
        "差 0.000009（阈值 0.0005）。",
        "本地与集群结果可以并进同一张表。",
    ], size=12, colour=INK, space_after=3)

    conclusion(slide, "治理约束：11-23-45 只用来选 epoch，16-08-46 只在定稿时评一次；本报告所有对外数字都带 test 列。")


def all_arms(prs):
    """Every arm that exists on the test set, caption and no-caption side by side."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "ALL ARMS ON THE HELD-OUT TEST SET",
           "全部路线 × 有无 caption：test 上的 AbsRel 一览")
    rows = [
        ["线", "输入", "GT 视角", "无 caption", "cap 训练 / 空 prompt", "cap 训练 / 真 caption"],
        ["a", "RGB", ("rgb", {"colour": BAD}), ("0.0844", {"bold": True}), "未训练", "未训练"],
        ["b", "Thermal", "thr",
         ("0.0884  ★ 热像最好", {"bold": True, "colour": GOOD}),
         ("0.0869  ★ 全表最好", {"bold": True, "colour": GOOD}), "0.0882"],
        ["c1", "Thermal", "thr", "0.1217", "未训练", "未训练"],
        ["c2", "Thermal", "thr", "0.1013", "0.1030", "0.1022"],
        ["d1", "Thermal", "thr", "0.0973", "未训练", "未训练"],
        ["d2", "Thermal", "thr", "0.0903", "0.0876", "0.0881"],
    ]
    align = [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT,
             PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT]
    table(slide, rows, MARGIN, 1.62, BODY_W, 2.70, [0.7, 1.3, 1.0, 2.2, 2.6, 2.6],
          size=11.5, align=align)

    write(textbox(slide, MARGIN, 4.56, BODY_W, 1.90), [
        ("★ 不带 caption 的最好是 b 的 0.0884；带 caption 的最好是 b 的 0.0869"
         "（用 caption 训练、推理给空 prompt），也是全项目最好的热像成绩。", {"bold": True}),
        ("⚠️ a 的 0.0844 不参与这个比较。", {"bold": True, "colour": BAD}),
        "它是 RGB 输入、按 RGB 视角 GT 评分（train_route_suite.py:439），"
        "与热像各线不是同一份真值，混在一起比大小是无效的。它在这里是「热像离 RGB 还有多远」的参照。",
        "「未训练」= 该线没有 caption 臂：a + caption 是第一周任务 4 的缺口（约 67 小时）；"
        "c1 与 d1 从未训过 caption 臂。所以 caption 的结论只覆盖 b / c2 / d2 三条。",
        "同一行内左右比较即为 caption 的效应；三条线的分解见后两页。"
        "c2 是唯一一条加 caption 反而变差的（0.1013 → 0.1030）。",
    ], size=11, colour=INK, space_after=3)
    conclusion(slide, "带 caption 与不带 caption 的最好成绩都在 b 线上；caption 只买到 0.0884 → 0.0869，"
                      "而且要靠推理时不喂 caption 才拿得到。")


def mechanism(prs):
    """One mechanism; the rest of the deck is its consequences.

    The findings were being presented as a list of unrelated results, which is
    unpresentable. They are not unrelated: each follows from where the
    supervision is and what was left unfrozen.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "ONE MECHANISM", "一个机制，其余都是它的推论")

    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(MARGIN), Inches(1.56),
                                 Inches(BODY_W), Inches(0.82))
    bar.fill.solid()
    bar.fill.fore_color.rgb = INK
    bar.line.fill.background()
    bar.shadow.inherit = False
    frame = bar.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = frame.margin_right = Inches(0.22)
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = ("稀疏 LiDAR 只监督到 26.5% 的像素，且集中在近中景。"
                "监督覆盖不到的地方，解冻了什么，就在那里漂。")
    run.font.name, run.font.size, run.font.bold = FONT, Pt(15), True
    run.font.color.rgb = WHITE

    rows = [
        ["现象", "为什么", "证据"],
        ["天空被判成近处",
         "天空一个 LiDAR 回波都没有 ＝ 完全无监督。解冻的 U-Net 在那里自由漂移。",
         "上带 <10 m 占比：c1 19% / b 33% / c2·d2 44%"],
        ["c1 天空最好、整体最差",
         "同一个原因的两面：它几乎没动 Lotus，所以既没被带偏，也没学会任务。",
         "上带 16.1 m（真值 14.8）；test 0.1217，六线最差"],
        ["condition 只在有监督处帮忙",
         "AnyThermal 特征全面优于 VAE latent（c1 vs d1 受控，九分层全胜）；"
         "但改善最小的正是 row/top，无 GT 的天空判据上 c1 反而更好。",
         "近景 −31.6% / 远端 −23.2% / row-top 仅 −16.1%"],
        ["caption 对远端几乎为零",
         "文本只经 cross-attention 进来，不携带几何。监督缺失的地方它补不上。",
         "b far 0.1680 → b+cap 0.1679"],
        [("原版 AnyThermal 反过来验证", {"colour": GOOD}),
         "它冻结骨干、只训一个头，先验全程保住 —— 正是「不解冻就不漂」。",
         "far/all ×1.13，我方最好 ×1.69"],
    ]
    table(slide, rows, MARGIN, 2.62, BODY_W, 3.20, [2.3, 5.6, 4.2], size=11)

    write(textbox(slide, MARGIN, 6.00, BODY_W, 0.55), [
        "所以「整体最好（b）／远端最好（d2）／天空最好（c1）」不是三个矛盾，"
        "而是同一个取舍的三个位置：冻多少。",
    ], size=11.5, colour=INK, bold=True)
    conclusion(slide, "这是一个可预测的失效模式，不是若干个 bug —— 后面每一页都是它的一个侧面。")


def anythermal_losses(prs):
    """Why the external model does not have our sky problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "WHY ANYTHERMAL HAS NO SKY PROBLEM",
           "AnyThermal 的两个损失：天空为什么没被带偏")
    rows = [
        ["", "阶段一：表征蒸馏", "阶段二：深度头"],
        ["损失", "CLS token 对比损失，把热像编码器对齐到 DINOv2 的 RGB 表征",
         "MiDaS：SSI 数据项 + 多尺度梯度匹配项"],
        ["要标注吗", ("不要 —— 自监督", {"bold": True, "colour": GOOD}),
         "要 —— MS2 train split 的稀疏 LiDAR"],
        ["监督覆盖", ("整张图，天空在内", {"bold": True, "colour": GOOD}),
         ("26.5% 的像素，天空为零", {"colour": BAD})],
        ["骨干状态", "训练中", ("冻结，只训头", {"bold": True, "colour": GOOD})],
    ]
    table(slide, rows, MARGIN, 1.62, BODY_W, 1.95, [1.3, 5.6, 5.2], size=11)

    write(textbox(slide, MARGIN, 3.86, 6.05, 2.30), [
        ("为什么天空没问题", {"bold": True, "size": 13}),
        "天空的「知识」是阶段一学的，而阶段一的损失根本不看深度标注 —— "
        "对它来说天空不是无监督区。",
        "阶段二虽然只有 26.5% 的监督，但骨干已经冻住，它动不了那份先验，所以漂不起来。",
        "合起来是一个顺序：先用不需要标注的目标把先验建好，再冻住，"
        "最后才让稀疏标注去调一个小头。",
    ], size=11, colour=INK, space_after=3)

    write(textbox(slide, 6.95, 3.86, 5.78, 2.30), [
        ("对照我方", {"bold": True, "size": 13}),
        "我们没有阶段一。先验来自 Stable Diffusion 的预训练，"
        "但阶段二把那 867 M 全解冻了 —— 等于把唯一的先验交给 26.5% 的监督去改写。",
        ("§8 那个失败现在有了解释：", {"bold": True, "colour": BAD}),
        "我们借的是它的第二个损失项（多尺度梯度匹配），三条臂单调恶化。"
        "回头看是借错了 —— 梯度匹配同样只在有效像素上算，天空照样没监督。"
        "关键不在第二个损失项，在第一个阶段，以及之后的冻结。",
    ], size=11, colour=INK, space_after=3)

    conclusion(slide, "可借鉴的不是那条损失，是「无标注目标建先验 → 冻住 → 稀疏标注只训小头」这个顺序。")


def main_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "BASELINE · SIX ROUTES ON THE HELD-OUT TEST SET", "主表：六条路线的 val 与 test")
    rows = [
        ["线", "输入", "Condition", "可训练参数", "val", "test", "rmse", "a1"],
        ["a  RGB+U-Net", "RGB", "冻结 VAE latent", "867.57 M", "0.0791", ("0.0844", {"bold": True}), "3.827", "0.9269"],
        ["b  Thermal+U-Net", "Thermal", "冻结 VAE latent", "867.57 M", ("0.0775", {"bold": True}), "0.0884", "3.956", "0.9071"],
        ["c1  VAE-adapter", "Thermal", "VAE latent + Adapter", "7.11 M", "0.1127", "0.1217", "5.465", "0.8365"],
        ["c2  VAE-adapter→U-Net", "Thermal", "VAE latent + 冻结 Adapter", "867.57 M", "0.0932", "0.1013", "4.418", "0.8826"],
        ["d1  AnyTh-adapter", "Thermal", "AnyThermal 特征 + Adapter", "9.41 M", "0.0860", "0.0973", "4.432", "0.8883"],
        ["d2  AnyTh-adapter+U-Net", "Thermal", "AnyThermal 特征 + Adapter", "876.98 M", "0.0815", "0.0903", "3.855", "0.9086"],
    ]
    align = [PP_ALIGN.LEFT] + [PP_ALIGN.LEFT] * 2 + [PP_ALIGN.RIGHT] * 5
    table(slide, rows, MARGIN, 1.62, BODY_W, 3.10, [2.5, 1.0, 2.9, 1.5, 1.0, 1.0, 1.0, 1.0],
          size=11.5, align=align)

    write(textbox(slide, MARGIN, 4.92, BODY_W, 1.30), [
        "· 热像几乎追平 RGB：a 0.0844 vs b 0.0884，仅差 4.5%。但 b 的跨序列泛化更差 —— "
        "val→test 退化 +0.0109，a 只有 +0.0053。",
        "· d 系全面压过 c 系（d1 0.0973 vs c1 0.1217；d2 0.0903 vs c2 0.1013）："
        "AnyThermal 特征作为 condition 优于 VAE latent，在两个参数量级上都成立。",
        "· d1 用 9.41 M（1%）参数达到 0.0973，逼近 867 M 的 b。",
    ], size=12, colour=INK, space_after=4)
    footnote(slide, "外部参照点（原版 AnyThermal，test 0.0821）在任务 4 那一页有完整分层表，此处不重复。")
    conclusion(slide, "六线的排序稳定：condition 用什么，比在 condition 上挂多少参数更要紧。")


IN_RGB = {"kind": "image", "file": "in_rgb.png", "w": 1.30,
          "label": "RGB x", "sub": "3 × 256 × 640"}
IN_THR = {"kind": "image", "file": "in_thermal.png", "w": 1.30,
          "label": "热像 x", "sub": "3 × 256 × 640"}
def out_disp(key):
    """That route's own prediction on the same frame as the input panel.

    Both come from frame 002403 of the test sequence: the input straight from
    sync_data, the output from export_qualitative.py's `<label>_pred_demo.png`,
    which is `colorize_depth_map(pred, reverse_color=True)` on the last frame it
    rendered.  Pairing an RGB input with some other scene's thermal output, as
    the first draft did, makes the figure a lie about what the route consumes.
    """
    return {"kind": "image", "file": f"out_{key}.png", "w": 1.30,
            "label": "视差 ŷ", "sub": "256 × 640"}
VAE_ENC = {"kind": "enc", "w": 1.05, "label": "Variational Encoder", "sub": "34.16 M"}
VAE_DEC = {"kind": "dec", "w": 1.05, "label": "Variational Decoder", "sub": "49.49 M · fp32"}
CONCAT = {"kind": "concat", "w": 0.44}
Z_X = {"kind": "cube", "w": 0.58, "h": 0.68, "label": "z_x", "sub": "4 × 32 × 80"}
Z_C = {"kind": "cube", "w": 0.58, "h": 0.68, "label": "z_c", "sub": "4 × 32 × 80"}
X0 = {"kind": "cube", "w": 0.58, "h": 0.68, "label": "x̂₀", "sub": "4 × 32 × 80"}
FEAT = {"kind": "cube", "w": 0.58, "h": 0.68, "label": "特征金字塔", "sub": "ViT-B/14 tokens"}
ANY_ENC = {"kind": "enc", "w": 1.05, "label": "AnyThermal Encoder", "sub": "86.58 M · DINOv2",
           "fill": ANY_FILL, "line": ANY_LINE}


def unet(trained):
    return {"kind": "unet", "w": 1.70, "h": 1.04, "trained": trained,
            "label": "Lotus U-Net", "sub": "867.57 M · 单步"}


def adapter(kind, trained=True):
    if kind == "vae":
        return {"kind": "block", "w": 0.90, "h": 1.00, "trained": trained,
                "label": "Adapter", "sub": "7.11 M · 残差"}
    return {"kind": "block", "w": 0.90, "h": 1.00, "trained": trained,
            "label": "Adapter", "sub": "9.41 M · V2.3"}


ROUTES = [
    {"key": "a", "kicker": "ROUTE A · DATA FLOW",
     "title": "a 线：RGB → 冻结 VAE → 可训练 U-Net",
     "nodes": [IN_RGB, VAE_ENC, Z_X, CONCAT, unet(True), X0, VAE_DEC, None],
     "note": "与 b 线逐位相同，唯一变量是输入模态。它是「热像离 RGB 还有多远」的对照，"
             "也是唯一按 RGB 视角 GT 评分的线。",
     "conclusion": "唯一拿梯度的是 867.57 M 的 U-Net；VAE 与文本通路全部冻结。test 0.0844。"},
    {"key": "b", "kicker": "ROUTE B · DATA FLOW",
     "title": "b 线：Thermal → 冻结 VAE → 可训练 U-Net",
     "nodes": [IN_THR, VAE_ENC, Z_X, CONCAT, unet(True), X0, VAE_DEC, None],
     "note": "热像被当成三通道图直接喂进为 RGB 训练的 VAE。这条线是本项目最好的热像成绩"
             "（val 0.0775），也是其余各线的比较基准。",
     "conclusion": "与 a 线同一套可训练参数；condition 换成热像的 VAE latent。test 0.0884。"},
    {"key": "c1", "kicker": "ROUTE C1 · DATA FLOW",
     "title": "c1 线：Thermal → 冻结 VAE → 可训练 Adapter → 冻结 U-Net",
     "nodes": [IN_THR, VAE_ENC, Z_X, adapter("vae"), Z_C, CONCAT, unet(False), X0,
               VAE_DEC, None],
     "note": "Adapter 是 latent 上的残差 CNN，零初始化 ⇒ 未训练时是恒等映射，"
             "所以起点恰好是「冻结直推」。全链路只有 0.5% 的参数拿梯度。",
     "conclusion": "只训 7.11 M 的 Adapter，U-Net 冻结。test 0.1217 —— 六线里最差。"},
    {"key": "c2", "kicker": "ROUTE C2 · DATA FLOW（任务 1b 阶段二）",
     "title": "c2 线：先训 Adapter（= c1），再冻结它、只训 U-Net",
     "nodes": [IN_THR, VAE_ENC, Z_X, adapter("vae", trained=False), Z_C, CONCAT,
               unet(True), X0, VAE_DEC, None],
     "note": "这是导师要的「先训 Adapter 再训后面的 U-Net」：阶段一就是已完成的 c1，"
             "阶段二用 --init-from 载入它、--freeze-adapter 冻住，只让 U-Net 拿梯度。"
             "所以可训练参数与 b 线相同，唯一差别是 condition 先过了一层 c1 训出来的 Adapter。",
     "conclusion": "只训 U-Net，867.57 M（Adapter 继承自 c1、冻结）。test 0.1013，远端 0.1817（全线最差）。"},
    {"key": "d1", "kicker": "ROUTE D1 · DATA FLOW",
     "title": "d1 线：Thermal → 冻结 AnyThermal → 可训练 Adapter → 冻结 U-Net",
     "nodes": [IN_THR, ANY_ENC, FEAT, adapter("any"), Z_C, CONCAT, unet(False), X0,
               VAE_DEC, None],
     "note": "注意 VAE encoder 整条不在了 —— condition 改由 AnyThermal 的 DINOv2 特征经 Adapter 生成。"
             "VAE 只剩解码端。",
     "conclusion": "9.41 M（全链路 0.7%）达到 test 0.0973，逼近 867 M 的 b 线。"},
    {"key": "d2", "kicker": "ROUTE D2 · DATA FLOW",
     "title": "d2 线：Thermal → 冻结 AnyThermal → 可训练 Adapter + 可训练 U-Net",
     "nodes": [IN_THR, ANY_ENC, FEAT, adapter("any"), Z_C, CONCAT, unet(True), X0,
               VAE_DEC, None],
     "note": "本轮远端表现最好的线。§7 的 donor-swap 证明这条特征通路真的在承载信号："
             "推理时把它换成随机 donor 帧的，误差从 0.0903 翻到 0.2971。",
     "conclusion": "876.98 M。test 0.0903，但远端 0.1524 —— 六线里唯一在远端改善的。"},
]


def legend(slide, y):
    """Shape says what a block is; colour says whether it takes gradients."""
    items = [
        ("enc", VAE_FILL, VAE_LINE, "VAE 编解码"),
        ("enc", ANY_FILL, ANY_LINE, "AnyThermal"),
        ("enc", TXT_FILL, TXT_LINE, "文本编码"),
        ("unet_hot", HOT_FILL, HOT_LINE, "拿梯度"),
        ("unet_cold", COLD_FILL, COLD_LINE, "冻结"),
        ("cube", CUBE_FILL, CUBE_LINE, "张量"),
    ]
    x = MARGIN + 0.30
    for kind, fill, line, label in items:
        if kind == "cube":
            _shape(slide, MSO_SHAPE.CUBE, x, y, 0.26, 0.30, fill=fill, line=line, lw=1.0)
        elif kind == "enc":
            _shape(slide, MSO_SHAPE.TRAPEZOID, x, y, 0.30, 0.34, rot=90, fill=fill, line=line, lw=1.0)
        else:
            _shape(slide, MSO_SHAPE.TRAPEZOID, x - 0.08, y, 0.16, 0.34, rot=90, fill=fill, line=line, lw=1.0)
            _shape(slide, MSO_SHAPE.TRAPEZOID, x + 0.08, y, 0.16, 0.34, rot=270, fill=fill, line=line, lw=1.0)
        write(textbox(slide, x + 0.22, y - 0.11, 1.5, 0.24), [label], size=9.5, colour=GREY)
        x += 1.62


def route_flow(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, spec["kicker"], spec["title"])
    # None marks the output panel: it is per-route, so it is resolved here and
    # every route shows its own prediction on the frame its input panel shows
    nodes = [n if n is not None else out_disp(spec["key"]) for n in spec["nodes"]]
    flow(slide, nodes)
    legend(slide, 5.86)
    write(textbox(slide, MARGIN, 6.14, BODY_W, 0.40), [spec["note"]], size=9.5, colour=GREY)
    conclusion(slide, spec["conclusion"])


def route_flows(prs):
    for spec in ROUTES:
        route_flow(prs, spec)


def caption_flow(prs):
    """The caption arm as a picture: same image path, only the string changes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "CAPTION ARM · WHERE THE TWO EFFECTS ENTER",
           "caption 是怎么进去的，以及两个效应各自在哪一步")
    nodes = [IN_THR, VAE_ENC, Z_X, CONCAT, unet(True), X0, VAE_DEC, out_disp("b")]
    flow(slide, nodes, prompts=['c = ""（空 prompt）', "c = 真 caption",
                                "c = 打乱的 caption"])

    write(textbox(slide, MARGIN, 5.40, BODY_W, 1.15), [
        ("① 权重效应 —— 训练时用哪种 prompt。", {"bold": True}),
        "改变的是 U-Net 学出来的权重，比较的是两次独立训练的 checkpoint，"
        "所以带 seed 不确定性（本轮单 seed）。",
        ("② 注入效应 —— 推理时喂哪种 prompt。", {"bold": True}),
        "同一份权重、同一批帧，除最左边那三个框外图上没有任何东西改变，所以没有混淆变量。",
        "caption 由 InternVL3-8B 从对应的 RGB 帧生成；「打乱」= 换成别帧的 caption，"
        "保留文本分布、只打断图文对应，用来区分「起作用的是内容」还是「有文本就行」。",
    ], size=10.5, colour=INK, space_after=2)
    conclusion(slide, "图像通路完全不动；两个效应的区别只在于「换的是权重」还是「换的是 prompt」。")


def caption_effects(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "★ CAPTION · TWO OPPOSITE EFFECTS", "caption 拆成两个效应后，两个都显著")
    write(textbox(slide, MARGIN, 1.48, BODY_W, 0.34), [
        "同一条线的三个考法（test 2,543 帧，AbsRel，越小越好）：",
    ], size=12, colour=GREY)
    rows = [
        ["线", "无 cap 训练", "cap 训练 / 空 prompt", "cap 训练 / 真 caption"],
        ["b  Thermal+U-Net", "0.0884", ("0.0869", {"bold": True, "colour": GOOD}), "0.0882"],
        ["c2  VAE-adapter+U-Net", "0.1013", "0.1030", "0.1022"],
        ["d2  AnyTh-adapter+U-Net", "0.0903", ("0.0876", {"bold": True, "colour": GOOD}), "0.0881"],
    ]
    align = [PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 3
    table(slide, rows, MARGIN, 1.86, BODY_W, 1.20, [3.2, 2.6, 3.2, 3.2], size=12, align=align)

    write(textbox(slide, MARGIN, 3.30, BODY_W, 0.34), [
        "拆成两个效应（负 = 更好；n=2,543 配对，bootstrap 3000 次）：",
    ], size=12, colour=GREY)
    rows = [
        ["线", "① 权重效应（训练时用 caption）", "② 注入效应（推理时喂 caption）"],
        ["b", ("−0.00153  [−0.00188, −0.00118]  显著", {"colour": GOOD}),
         ("+0.00130  [+0.00110, +0.00152]  显著", {"colour": BAD})],
        ["c2", ("+0.00169  [+0.00146, +0.00194]  显著", {"colour": BAD}),
         ("−0.00082  [−0.00104, −0.00062]  显著", {"colour": GOOD})],
        ["d2", ("−0.00270  [−0.00332, −0.00206]  显著", {"colour": GOOD}),
         ("+0.00046  [+0.00037, +0.00056]  显著", {"colour": BAD})],
    ]
    table(slide, rows, MARGIN, 3.68, BODY_W, 1.20, [1.6, 5.3, 5.3], size=12,
          align=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT])

    write(textbox(slide, MARGIN, 5.12, BODY_W, 1.10), [
        "· 全项目最好的热像成绩 0.0869 —— b 用 caption 训练、推理时给空 prompt。"
        "caption 的价值在训练期，不在推理期。",
        "· 注入税的大小由 condition 架构决定：b +0.00130，d2 只有 +0.00046，小 3 倍。"
        "语义层级的 condition 不能让注入变得有益，但能把代价降到三分之一。",
        "· c2 两个方向都反着来，与它在所有其他指标上的劣势一致 —— VAE latent 这条路整体是错的。",
    ], size=12, colour=INK, space_after=3)
    conclusion(slide, "以往「caption 不显著」是两个相反效应互相抵消的假象；拆开后两个都显著。")


def caption_strength(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "★ CAPTION · HOW STRONG IS EACH CLAIM", "两个效应的证据强度不同，必须分开说")
    rows = [
        ["", "② 注入效应", "① 权重效应"],
        ["比的是什么", "同一份权重、同一批帧，只有 prompt 不同",
         "两次独立训练的 checkpoint"],
        ["混淆变量", "没有 —— 除 prompt 外一切固定", "训练随机性（单 seed）"],
        ["能说到多硬", "无懈可击，可以直接下结论",
         "有力但不等于定论，汇报时标注「单 seed」"],
        ["旁证", "—",
         "三条线方向不一致（c2 反向）。若只是「重训一次就变好」的假象，三条应同向"],
    ]
    table(slide, rows, MARGIN, 1.70, BODY_W, 2.60, [1.9, 4.6, 5.6], size=12)

    write(textbox(slide, MARGIN, 4.55, BODY_W, 1.55), [
        ("为什么这一条值得单独一页", {"bold": True, "size": 13}),
        "① 和 ② 是同一张表里读出来的，但它们的可信度差一个量级。把它们并列成"
        "「caption 有两个效应」而不区分强度，是这轮最容易被问倒的地方。",
        "补齐 ① 的办法很明确：d2+caption 跑第二个 seed（约 27 小时），"
        "把单次观测升级成可复现观测。这是下一步的第一优先项。",
    ], size=12, colour=INK, space_after=5)
    conclusion(slide, "② 可以当结论讲；① 只能当「单 seed 下的强观测」讲，第二 seed 是明确的补救路径。")


def far_field(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "FAR FIELD · A CONTROLLED ATTRIBUTION", "d2 的远端优势，以及它能归因到什么")
    rows = [
        ["分层", "b", "b+cap", "d1", "c2", "d2"],
        ["all", "0.0884", "0.0882", "0.0973", "0.1013", "0.0903"],
        ["depth/far >30m", "0.1680", "0.1679", "0.1726", ("0.1817", {"colour": BAD}),
         ("0.1524", {"bold": True, "colour": GOOD})],
        ["row/top", "0.1572", "0.1598", "0.1767", "0.1858", "0.1705"],
        ["structure/boundary", "0.1075", "0.1089", "0.1184", "0.1242", "0.1112"],
        ["structure/interior", "0.0862", "0.0857", "0.0946", "0.0989", "0.0878"],
        ["锐利度 b/i", "1.25", "1.27", "1.25", "1.26", "1.27"],
        ["远端退化 far/all", "×1.90", "×1.90", "×1.77", "×1.79",
         ("×1.69", {"bold": True, "colour": GOOD})],
    ]
    align = [PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 5
    table(slide, rows, MARGIN, 1.62, 7.55, 3.55, [2.3, 1.0, 1.0, 1.0, 1.0, 1.0],
          size=11.5, align=align)

    write(textbox(slide, 8.40, 1.62, 4.33, 3.60), [
        ("受控归因：c1 vs d1", {"bold": True, "size": 13}),
        "两条线都只训 Adapter、U-Net 冻结，唯一差别是 condition 来自 VAE latent 还是 "
        "AnyThermal 特征 —— 没有第二个变量。",
        ("九个分层全部 d1 胜，全部显著，胜率 69–85%。", {"bold": True, "colour": GOOD}),
        ("但这是「全面更好」，不是「专治远端」：", {"bold": True, "colour": BAD}),
        "相对改善 近景 31.6% > 边界 24.8% > 远端 23.2% > 整体 20.0%，"
        "而 row/top 只有 16.1%，是九项里最低之一。",
        "",
        "d2 vs b 才是拿整体换远端的那一对：d2 整体略差（0.0903 vs 0.0884），"
        "但远端 0.1524 vs 0.1680（−9.3%），是六线里唯一在远端改善的。",
        "",
        ("任务 1b 的答案：没修好。", {"bold": True, "colour": BAD}),
        "c2（先训 Adapter 再训 U-Net）远端 0.1817、row/top 0.1858，都是所有线里最差。",
    ], size=11, colour=INK, space_after=4)

    footnote(slide,
             "⚠️ 分层只统计有 GT 的像素。MS2 每帧 valid 平均 26.5%，天空没有 LiDAR 回波 —— "
             "row/top 是「上三分之一里有回波的像素」（楼顶、电线杆），不是天空。措辞不能混。",
             top=5.42)
    conclusion(slide, "远端改善来自 condition 的语义层级，不是参数量、不是 adapter 结构 —— 这是一个受控结论。")


def donor_swap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "FALSIFICATION · DOES THE FEATURE BRANCH DO ANYTHING",
           "证伪对照：d2 真的在用 AnyThermal 特征吗")
    write(textbox(slide, MARGIN, 1.48, BODY_W, 0.70), [
        "担心的是：adapter 可能学会忽略 AnyThermal 特征，d2 退化成一个带额外参数的 b。"
        "损失曲线看不出这件事。做法是推理时把 condition 换成随机 donor 帧的（uniform，自配对已换走）。",
    ], size=12.5, colour=INK)
    rows = [
        ["", "abs_rel", "rmse", "a1"],
        ["d2 原始", "0.0903", "3.855", "0.9086"],
        ["d2 只打乱 AnyThermal 特征", ("0.2971", {"bold": True, "colour": BAD}), "10.005",
         ("0.5186", {"bold": True, "colour": BAD})],
        ["d2 打乱整个 condition", "0.2981", "10.034", "0.5171"],
        ["b 打乱整个 condition（标尺）", "0.2945", "9.952", "0.5206"],
    ]
    align = [PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 3
    table(slide, rows, MARGIN, 2.36, 7.30, 1.75, [3.4, 1.4, 1.4, 1.4], size=12, align=align)

    write(textbox(slide, 8.20, 2.36, 4.53, 3.30), [
        ("读法", {"bold": True, "size": 13}),
        "只换特征金字塔，误差翻 3.3 倍、a1 从 0.91 塌到 0.52。",
        "「只打乱特征」与「打乱全部」几乎相同（0.2971 vs 0.2981）——"
        "AnyThermal 特征承载了几乎全部 condition 信号，"
        "adapter 的 thermal 张量输入不独立贡献。",
        "崩塌幅度与 b 打乱 VAE latent 相当（0.2945）。",
    ], size=11.5, colour=INK, space_after=5)

    write(textbox(slide, MARGIN, 4.42, 7.30, 0.60), [
        "anythermal 模式只换特征、保留正确的热像张量，"
        "所以零结果不能推给「adapter 丢了图像」—— 这是该检验的设计关键。",
    ], size=11.5, colour=GREY)

    write(textbox(slide, MARGIN, 5.16, BODY_W, 1.00), [
        ("反过来的一半：adapter 也是必需的", {"bold": True, "size": 12.5}),
        "把 AnyThermal 特征经零参数 bridge 直接塞给冻结的 Lotus-G（不经 adapter、不训练），"
        "test AbsRel = 0.4697 —— 比上表「打乱 condition」的 0.2971 还差。"
        "特征金字塔的格式是冻结 U-Net 从没见过的，「乱但同分布」好过「不同分布」。",
    ], size=11.5, colour=INK, space_after=3)
    conclusion(slide, "两个方向都验了：adapter 在用 AnyThermal 特征（打乱就塌），而少了 adapter 这些特征也用不了。")


def external_reference(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "TASK 4 · THE EXTERNAL REFERENCE POINT", "原版 AnyThermal：它没有远端问题")
    rows = [
        ["", "all", "far >30m", "row/top", "boundary", "interior", "b/i", "far/all"],
        ["b（我方最佳 baseline）", "0.0884", "0.1680", "0.1572", "0.1075", "0.0862", "1.25",
         ("×1.90", {"colour": BAD, "bold": True})],
        ["原版 AnyThermal", ("0.0821", {"bold": True}), ("0.0930", {"bold": True, "colour": GOOD}),
         "0.1257", "0.0973", "0.0808", "1.20",
         ("×1.13", {"colour": GOOD, "bold": True})],
    ]
    align = [PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 7
    table(slide, rows, MARGIN, 1.62, BODY_W, 1.05, [2.6, 1.2, 1.4, 1.2, 1.3, 1.2, 0.9, 1.2],
          size=11.5, align=align)
    footnote(slide, "AnyThermal 用 ssi 对齐（独立作业），两表分层定义共享，可并排读。"
                    "复现验证：论文 Table IV 报 AbsRel 0.0883（day/night/rainy 平均），我方实测 0.0821（仅白天 test），量级吻合。",
             top=2.76)

    write(textbox(slide, MARGIN, 3.34, BODY_W, 0.34),
          ["它用了什么监督（论文 §III-B / §III-F）："], size=12.5, colour=INK, bold=True)
    rows = [
        ["", "原版 AnyThermal", "我方六线"],
        ["深度监督", "MS2 自己的 train split + 稀疏 LiDAR GT", "同左"],
        ["骨干", "DINOv2 ViT-B/14，CLS token 对比损失蒸馏（自监督、无标注）", "Stable Diffusion U-Net"],
        ["深度头", "MiDaS 架构，「其余部分不变」", "扩散单步 x0 预测"],
        ["训练时骨干", ("冻结，只训头", {"bold": True}), ("整个 U-Net 解冻微调", {"bold": True})],
    ]
    table(slide, rows, MARGIN, 3.72, BODY_W, 1.65, [1.7, 6.6, 3.8], size=11.5)

    write(textbox(slide, MARGIN, 5.55, BODY_W, 0.90), [
        ("机制假设：", {"bold": True}),
        "稀疏 LiDAR 在远端/天空无信号；全解冻的 U-Net 在无监督区自由漂移，"
        "而冻结骨干保住了 DINOv2 的语义先验。d2 部分借到了这个好处（用 AnyThermal 特征），"
        "但下游 U-Net 仍全解冻，所以只到 ×1.69。",
    ], size=11.5, colour=INK, space_after=2)
    conclusion(slide, "同一份稀疏 LiDAR 监督，差别在架构 —— 远端退化是我方方法特有的，不是热像深度的通病。")


def sky_band(prs, figure: Path | None):
    """The sky failure, and the one line that does not have it.

    This replaces an earlier three-route version of the same figure: c1 is the
    whole point and it was not in it.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "SKY · WHAT UNFREEZING THE U-NET COSTS",
           "天空判近是解冻 U-Net 的代价，不是 condition 的问题")
    picture_or_placeholder(
        slide, figure, MARGIN, 1.40, BODY_W, 2.90,
        ["【待填：五线对比图】"])

    rows = [
        ["线", "Adapter", "U-Net", "上带预测中位", "test AbsRel"],
        ["c1", "训练", ("冻结", {"bold": True}),
         ("16.1 m", {"bold": True, "colour": GOOD}), ("0.1217", {"colour": BAD})],
        ["d1", "训练", ("冻结", {"bold": True}), "12.1 m", "0.0973"],
        ["b", "—", "训练", "12.0 m", "0.0884"],
        ["d2", "训练", "训练", "10.7 m", ("0.0903", {"bold": True})],
        ["c2", "冻结（继承 c1）", "训练", ("10.3 m", {"colour": BAD}), "0.1013"],
        [("真值（有回波像素）", {"colour": GREY}), "—", "—",
         ("14.8 m", {"bold": True}), "—"],
    ]
    align = [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT]
    table(slide, rows, MARGIN, 4.40, 6.60, 2.00, [0.7, 1.9, 1.0, 1.5, 1.2],
          size=10, align=align)

    write(textbox(slide, 7.45, 4.40, 5.28, 2.05), [
        ("c1 的上带最远，但不是没有红", {"bold": True, "size": 12}),
        "五帧的上带中位全部最高（14.7 / 15.2 / 18.1 / 18.9 / 16.1），也是唯一一条中位"
        "高于真值中位的；解冻 U-Net 的三条都掉到 10–12 m。"
        "但 c1 的上带仍有 6–32% 的像素在 10 m 内（帧 000703 是 32%，与 b 的 40% 差得不远）。",
        ("两个代价：", {"bold": True, "colour": BAD}),
        "① c1 的整体精度是六线最差（0.1217）—— 它保住 Lotus 的先验，正因为几乎没动 Lotus。"
        "② 上三分之一不全是天空：000703 是林荫路，上部的树枝本来就该判近，"
        "所以那 32% 里有一部分是对的。这个判据分不出「天空判错」和「近物判对」。",
        ("上带中位 = np.median(aligned[上三分之一])，不带 valid 掩码，天空算在内；"
         "真值行只统计有回波像素。五帧，按远端挑的。", {"size": 9, "colour": GREY}),
    ], size=10, colour=INK, space_after=3)

    conclusion(slide, "稀疏 LiDAR 在天空没有监督：解冻的 U-Net 在那里自由漂移，冻结的保住了预训练先验 —— "
                      "与原版 AnyThermal 冻结骨干、远端只退化 ×1.13 同一个机制。")


def caption_visual(prs, strips: dict[str, Path]):
    """The caption effect is real and invisible; say both.

    One strip per route rather than one nine-column block: at eleven columns the
    header of each panel overran its neighbour and nothing was readable.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "CAPTION · SIGNIFICANT AND INVISIBLE",
           "caption 的效应显著，但肉眼看不见")

    strip_w = 7.55
    for index, key in enumerate(("b", "c2", "d2")):
        top = 1.50 + index * 1.62
        picture_or_placeholder(slide, strips.get(key), MARGIN, top, strip_w, 1.44,
                               [f"【待填：{key} 线的 caption 对比图】"])

    write(textbox(slide, 8.42, 1.50, 4.31, 4.70), [
        ("三条线，每条三种 prompt", {"bold": True, "size": 12.5}),
        "上到下：b / c2 / d2。每条线的三格是同一份权重、同一批帧，"
        "只有喂进 CLIP 的字符串不同：空 prompt、真 caption、打乱的 caption。",
        ("看不出差别 —— 这正是该有的样子。", {"bold": True}),
        "换 prompt 让上带深度中位最多只动 0.54 m（9–14 m 的量级上，约 3–5%），"
        "多数不到 0.3 m。栏首那三个 median 常常连整数都一样。",
        "注入效应在 test 上是 b +0.00130 / c2 −0.00082 / d2 +0.00046，"
        "而整体 AbsRel 是 0.09–0.10 —— 千分之一二的量级，统计上显著、方向明确，但小到画不出来。",
        ("这一页的作用是防止过度解读：", {"bold": True, "colour": BAD}),
        "「caption 有帮助」在图上找不到支撑，能支撑它的只有配对检验的区间。",
        ("第三格「打乱的 caption」是关键对照：", {"size": 10.5}),
        ("它保留文本分布、只打断图文对应。与第二格无差别 ⇒ 起作用的是「有文本」而非内容。",
         {"size": 10.5, "colour": GREY}),
    ], size=10.5, colour=INK, space_after=3)
    conclusion(slide, "显著 ≠ 可见。caption 的效应要用配对区间讲，不能用图讲。")


def rain(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "SUPPLEMENT · WHEN DOES INJECTION HELP", "补充一：雨天 —— 注入效应的条件边界")
    write(textbox(slide, MARGIN, 1.46, BODY_W, 0.34), [
        "三个考卷，同一份 b+caption 权重，只有 prompt 内容不同（②注入效应口径）。",
    ], size=12, colour=GREY)
    rows = [
        ["考卷", "帧数", "b 的 abs_rel", "注入效应", "胜率", "caption 词表(等n)", "唯一率"],
        ["晴  16-08-46", "2,543", "0.0884", ("+0.00130  显著（有害）", {"colour": BAD}), "—", "810", "92.7%"],
        ["轻雨  16-19-00", "2,374", "0.0930",
         ("−0.00169  显著（有帮助）", {"colour": GOOD, "bold": True}), "60.7%",
         ("831", {"bold": True}), "90.9%"],
        ["重雨  16-59-13", "2,161", "0.1170", ("+0.00049  显著（有害）", {"colour": BAD}), "45.0%",
         ("709", {"bold": True}), "87.2%"],
    ]
    align = [PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.LEFT,
             PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT]
    table(slide, rows, MARGIN, 1.84, BODY_W, 1.35, [1.9, 0.9, 1.4, 3.3, 0.9, 1.9, 1.0],
          size=11.5, align=align)

    rows = [
        ["考卷", "视觉通路", "caption 信息量", "结果"],
        ["晴", "强", "高（810）", "文本冗余 → 注入税"],
        ["轻雨", ("被削弱", {"bold": True}), ("最高（831）", {"bold": True}),
         ("补位成功", {"bold": True, "colour": GOOD})],
        ["重雨", "严重削弱", "最低（709）", "文本自己也失效 → 注入税"],
    ]
    table(slide, rows, MARGIN, 3.44, 7.10, 1.35, [1.2, 1.6, 1.8, 3.0], size=11.5)

    write(textbox(slide, 8.00, 3.44, 4.73, 2.60), [
        ("两个条件必须同时满足", {"bold": True, "size": 12.5}),
        "视觉通路被削弱 且 文本仍携带有效信息。",
        "这修正了原先的「文本价值 = 视觉通路有多弱」—— "
        "那个说法在 caption 源与模型输入同模态时看不出区别；"
        "本设置 caption 来自 RGB、模型吃热像，两个变量解耦后才露出真正的驱动因素。",
    ], size=11.5, colour=INK, space_after=4)

    write(textbox(slide, MARGIN, 5.02, 7.10, 1.15), [
        ("未复现的部分：", {"bold": True, "colour": BAD}),
        "d2 上轻雨与重雨的注入效应几乎相同（−0.00059 / −0.00062），尽管词表差 15%。"
        "「效应跟着 caption 信息量走」在 b 上成立、在 d2 上不成立。",
    ], size=11, colour=INK, space_after=2)
    conclusion(slide, "可证伪的下一步：夜间 21-58-13（RGB 近全黑）。预测词表低于重雨、注入税大于 +0.0005。")


def gradient_matching(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "SUPPLEMENT · A CLEAN NEGATIVE RESULT",
           "补充二：MiDaS 多尺度梯度匹配损失 —— 试了，失败了")
    write(textbox(slide, MARGIN, 1.46, BODY_W, 0.80), [
        "动机：逐点 L1 让「边界抹平」比「锐利但错位」便宜（合成阶跃上 0.0385 vs 0.0729），"
        "而 AbsRel/RMSE 同为逐点平均，奖励同一种投机 —— 这解释了「数值好看但可视化模糊」。"
        "MiDaS 的损失有第二项（梯度匹配），我方只有数据项。",
    ], size=12, colour=INK)
    rows = [
        ["臂", "best epoch", "val abs_rel", "rmse", "a1"],
        ["w=0", "2", ("0.07974", {"bold": True, "colour": GOOD}), "3.464", "0.9256"],
        ["w=5", "3", "0.08539", "3.760", "0.9136"],
        ["w=20", "2", ("0.09284", {"colour": BAD}), "4.277", "0.9011"],
    ]
    align = [PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 4
    table(slide, rows, MARGIN, 2.46, 6.10, 1.35, [1.2, 1.4, 1.6, 1.2, 1.2], size=12, align=align)
    footnote(slide, "三条臂各 3 epoch，单调恶化。", top=3.92)

    write(textbox(slide, 7.05, 2.46, 5.68, 3.40), [
        ("失败机制（有剂量-反应关系、有机制，可写）", {"bold": True, "size": 12.5}),
        "alignment_scale 从 0.26 涨到 3.3 —— 输出动态范围被压缩 12 倍；"
        "clamped_above 从 5 涨到 80。",
        "scale 虽 detach，但出现在损失对预测的梯度里，形成正反馈 —— "
        "w=20 的梯度范数常态触顶 --max-grad-norm 1.0。",
        "",
        ("实现本身是对的：", {"bold": True}),
        "合成阶跃上正确锐利预测两项均为 0，模糊预测梯度项 0.0262；"
        "26.5% 稀疏 mask 下 0.026164 vs 稠密 0.026163 —— "
        "多尺度掩码在 MS2 的 GT 密度下成立。",
    ], size=11.5, colour=INK, space_after=4)

    write(textbox(slide, MARGIN, 4.35, 6.10, 1.70), [
        ("学到的方法论", {"bold": True, "size": 12.5}),
        "权重标定按损失值之比定是错的：合成数据比值 0.68，"
        "真实数据实测 gmatch/gt_ssi_l1 = 0.087，差 8 倍。",
        "应该按梯度量级定权重，不是按损失值。",
    ], size=11.5, colour=INK, space_after=4)
    conclusion(slide, "负面结果，但不是空手：机制清楚、有剂量-反应关系，且给出了一条可复用的方法论。")


def limits(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "LIMITS", "局限（这些必须写进结论）")
    rows = [
        ["", "局限", "影响到哪条结论"],
        ["1", "GT 覆盖率 26.5%，天空无 LiDAR 回波 —— AbsRel/RMSE/所有分层都测不到天空。"
              "RMSE 逃不出这个问题（同一个 valid mask），它只是单位可读。",
         "所有「天空」措辞。分层里的 row/top 是「上三分之一有回波的像素」，不是天空"],
        ["2", "权重效应（①）单 seed；注入效应（②）无此问题。", "caption 结论的前半句"],
        ["3", "a + caption 未跑（第一周任务 4 缺口，67 小时）。", "RGB 线上的 caption 效应无数据"],
        ["4", "对比可视化尚未回来（第二周任务 3 的后半）。", "可视化页"],
        ["5", "雨天两因素模型在 d2 上未复现。", "「效应跟着 caption 信息量走」只在 b 上成立"],
        ["6", "AbsRel/RMSE/所有分层都测不到天空。不依赖 GT 的判据已有"
              "（export_qualitative.py 的上带中位深度），但只在 5 帧上跑过，尚未全量。",
         "「它没有天空问题」这类表述只能说到「有回波的上三分之一」"],
    ]
    table(slide, rows, MARGIN, 1.66, BODY_W, 4.40, [0.4, 7.0, 4.0], size=11)
    conclusion(slide, "最大的一条是 GT 覆盖率：本轮所有关于天空的说法，严格讲都只覆盖「有回波的上三分之一」。")


def next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "NEXT", "下一步，按优先级")
    rows = [
        ["", "做什么", "为什么", "成本"],
        ["1", "在无监督区加约束（而不是换 condition）",
         "机制页的直接推论：冻结更多，或补一个不依赖 GT 的正则。"
         "最省的一步是把 c1 的「冻结 U-Net」和 d2 的「语义 condition」合起来试一条线",
         "一条线一次训练"],
        ["2", "d2 + caption 第二个 seed",
         "把权重效应 ① 从单次观测升级为可复现观测 —— 这是本轮最软的一环", "约 27 小时"],
        ["3", "上带深度判据跑全量",
         "工具已有（export_qualitative.py），现在只在 5 帧上跑过；换成全部 2,543 帧就能"
         "把「天空判成近处」从佐证变成结论 —— 这是局限 1 唯一能被绕开的部分", "一个 GPU 作业"],
        ["4", "夜间 21-58-13 第四个数据点",
         "雨天两因素模型的可证伪预测：词表低于重雨、注入税大于 +0.0005", "一个评估作业"],
        ["5", "a + caption",
         "补齐第一周任务 4 的缺口，看 caption 双效应在 RGB 线上是否同构", "约 67 小时"],
    ]
    table(slide, rows, MARGIN, 1.70, BODY_W, 3.10, [0.4, 3.3, 6.4, 1.7], size=11.5)
    write(textbox(slide, MARGIN, 5.10, BODY_W, 1.10), [
        "1 是机制页指出的方向，2 决定 caption 结论能讲多硬；3–5 是扩展。"
        "集群上的产物路径、自查工具与全部原始数字见 docs/AIRE_RESULTS_20260802.md。",
    ], size=12, colour=GREY)
    conclusion(slide, "先补 ① 的第二 seed 和零训练参照点，再谈扩展。")


COMPARE = Path(__file__).resolve().parents[1] / "docs" / "figures" / "compare"
FIG5 = COMPARE / "qual_5routes" / "comparison_strip_shared.png"
FIGCAP = {k: COMPARE / f"qual_cap_{k}" / "comparison_strip.png"
          for k in ("b", "c2", "d2")}


def build(output: Path, figure: Path | None) -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    cover(prs)
    overview(prs)
    protocol(prs)
    main_table(prs)
    all_arms(prs)
    mechanism(prs)
    route_flows(prs)
    caption_flow(prs)
    caption_visual(prs, FIGCAP)
    caption_effects(prs)
    caption_strength(prs)
    far_field(prs)
    donor_swap(prs)
    external_reference(prs)
    anythermal_losses(prs)
    sky_band(prs, figure or FIG5)
    rain(prs)
    gradient_matching(prs)
    limits(prs)
    next_steps(prs)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"{output}  ({len(prs.slides)} slides)")


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__,
                                     formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--figure", type=Path, default=None,
                        help="Comparison PNG from tools/export_qualitative.py. "
                             "Omitted -> the qualitative slide keeps its marked placeholder.")
    args = parser.parse_args()
    build(args.output, args.figure)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
