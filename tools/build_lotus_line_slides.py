"""Build the Lotus-line stage-review slide deck.

All qualitative images embedded by this script are untouched files produced by
the upstream Iris/Lotus ``colorize_depth_map`` visualizer.  The deck is for the
route-selection phase only; it does not mix in the final-paper ``ms2_eval``
protocol or its visualizations.
"""

from __future__ import annotations

import csv
import json
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / ".codex" / "vendor"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = ROOT / "docs" / "slides"
OUT_PATH = OUT_DIR / "lotus_line_stage_review_20260703_v2.pptx"

NAVY = RGBColor(13, 25, 45)
NAVY_2 = RGBColor(20, 38, 65)
BLUE = RGBColor(48, 108, 255)
CYAN = RGBColor(47, 197, 225)
GREEN = RGBColor(26, 177, 113)
YELLOW = RGBColor(255, 194, 71)
RED = RGBColor(239, 91, 91)
WHITE = RGBColor(246, 249, 255)
MUTED = RGBColor(165, 180, 204)
GRID = RGBColor(53, 74, 105)
FONT = "Microsoft YaHei"


def add_rect(slide, x, y, w, h, fill, radius=False, line=None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_text(slide, text, x, y, w, h, size=18, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rich_text(slide, parts, x, y, w, h, size=18, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = align
    for text, color, bold in parts:
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_title(slide, title, kicker, index):
    add_text(slide, kicker.upper(), 0.55, 0.28, 8.5, 0.26, 9, CYAN, True)
    add_text(slide, title, 0.55, 0.57, 11.8, 0.62, 25, WHITE, True)
    add_text(slide, f"{index:02d}", 12.15, 0.38, 0.55, 0.35, 11, MUTED, True,
             PP_ALIGN.RIGHT)
    add_rect(slide, 0.55, 1.23, 12.15, 0.018, GRID)


def add_footer(slide, text="Route selection · upstream Iris/Lotus evaluator"):
    add_text(slide, text, 0.55, 7.18, 12.1, 0.18, 8, MUTED)


def set_background(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY


def add_bullets(slide, items, x, y, w, h, size=16, accent=CYAN):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)
        p.level = 0
        p.text = "•  " + p.text
    add_rect(slide, x - 0.15, y + 0.04, 0.035, min(h - 0.08, 0.46), accent)
    return box


def add_metric_card(slide, x, y, w, h, label, value, note, color=BLUE):
    add_rect(slide, x, y, w, h, NAVY_2, True, GRID)
    add_text(slide, label, x + 0.18, y + 0.15, w - 0.36, 0.24, 10, MUTED, True)
    add_text(slide, value, x + 0.18, y + 0.48, w - 0.36, 0.48, 24, color, True)
    add_text(slide, note, x + 0.18, y + h - 0.43, w - 0.36, 0.27, 9, MUTED)


def add_table(slide, x, y, w, h, columns, rows, widths=None, font_size=11,
              highlight_cols=None, highlight_cells=None):
    table = slide.shapes.add_table(len(rows) + 1, len(columns), Inches(x), Inches(y),
                                   Inches(w), Inches(h)).table
    if widths:
        for col, width in zip(table.columns, widths):
            col.width = Inches(width)
    highlight_cols = highlight_cols or set()
    highlight_cells = highlight_cells or set()
    for c, col in enumerate(columns):
        cell = table.cell(0, c)
        cell.text = col
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE if c in highlight_cols else NAVY_2
    for r, row in enumerate(rows, 1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(27, 47, 78) if r % 2 else RGBColor(22, 41, 69)
            if (r - 1, c) in highlight_cells:
                cell.fill.fore_color.rgb = RGBColor(24, 91, 82)
    for r in range(len(rows) + 1):
        for c in range(len(columns)):
            cell = table.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.025)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font_size)
                    run.font.bold = r == 0 or c == 0
                    run.font.color.rgb = WHITE
    return table


def add_picture_contain(slide, path, x, y, w, h, border=GRID):
    from PIL import Image

    path = Path(path)
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    add_rect(slide, x, y, w, h, RGBColor(8, 15, 28), False, border)
    slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))


def route_vis_path(directory, seq, sample_id):
    return ROOT / "outputs" / "lotus_line_v1" / directory / "vis" / (
        f"sync_data__{seq}_thr_img_left_pred_{sample_id}.png"
    )


def metrics_for_sample(csv_path, sample_id):
    with open(csv_path, encoding="utf-8", newline="") as f:
        for row in csv.DictReader(f):
            if row["filename"].endswith(f"pred_{sample_id}"):
                return {k: float(v) for k, v in row.items() if k != "filename"}
    raise KeyError(sample_id)


def build():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    route_data = json.loads((ROOT / "outputs" / "lotus_line_v1" / "route_caption_analysis.json").read_text(encoding="utf-8"))
    test_data = json.loads((ROOT / "outputs" / "lotus_line_v1" / "unet_caption_test_paired_analysis.json").read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 — cover
    s = prs.slides.add_slide(blank); set_background(s)
    add_rect(s, 0, 0, 13.333, 7.5, NAVY)
    add_rect(s, 0.62, 0.67, 0.08, 1.05, BLUE)
    add_text(s, "LOTUS LINE · STAGE REVIEW", 0.86, 0.7, 6.5, 0.3, 10, CYAN, True)
    add_text(s, "从路线选择到 Caption 贡献验证", 0.86, 1.16, 11.5, 0.8, 30, WHITE, True)
    add_text(s, "MS2 left thermal · Iris/Lotus official evaluation · Val → frozen Test", 0.88, 2.12, 10.8, 0.34, 15, MUTED)
    add_metric_card(s, 0.86, 3.25, 2.8, 1.55, "VAL ROUTE SELECTION", "5,810", "同一序列 / 同一协议", BLUE)
    add_metric_card(s, 3.92, 3.25, 2.8, 1.55, "FROZEN TEST", "9,508", "独立序列 / 无调参", GREEN)
    add_metric_card(s, 6.98, 3.25, 2.8, 1.55, "PAIRED BOOTSTRAP", "2,000×", "Correct vs Empty", YELLOW)
    add_text(s, "阶段性结论：U-Net-only 是当前最干净的 Caption 验证路线；Caption 有稳定贡献，但并非所有指标全面提升。", 0.88, 5.35, 11.2, 0.75, 18, WHITE, True)
    add_footer(s, "Stage review · 2026-07-03")

    # 2 — protocol
    s = prs.slides.add_slide(blank); set_background(s); add_title(s, "先统一比较协议，再谈模型优劣", "Evaluation protocol", 2)
    add_bullets(s, [
        "输入统一：MS2 left thermal；GT 统一：thermal-view filtered LiDAR depth。",
        "四条路线均调用 lotus/evaluation/evaluation.py::evaluation_depth，采用 least_square_disparity 对齐。",
        "所有预测图均来自官方 colorize_depth_map；不使用自定义色图、误差拼图或逐图归一化。",
        "Val 用于路线选择；选定 checkpoint 后冻结，在独立 Test 序列上只做一次确认。",
    ], 0.78, 1.55, 7.15, 3.85, 16)
    add_metric_card(s, 8.38, 1.58, 3.95, 1.32, "MANIFEST SHA256 (VAL)", "c8b63d0a…7301e", "5,810 samples", BLUE)
    add_metric_card(s, 8.38, 3.08, 3.95, 1.32, "MANIFEST SHA256 (TEST)", "83b96c92…f60b", "9,508 samples", GREEN)
    add_metric_card(s, 8.38, 4.58, 3.95, 1.32, "THERMAL DECODE", "uint16 → uint8", "先解码，再复制为 3 channels", YELLOW)
    add_rich_text(s, [("边界：", YELLOW, True), ("本 deck 是路线选择结果，不等同于最终论文统一 MS2 benchmark。", WHITE, False)], 0.78, 6.25, 11.6, 0.35, 13)
    add_footer(s)

    # 3 — routes
    s = prs.slides.add_slide(blank); set_background(s); add_title(s, "四条路线回答四个不同问题", "Experimental design", 3)
    cards = [
        ("Direct", "0 参数桥接", "AnyThermal feature 直接映射到 Lotus-G condition", BLUE),
        ("Adapter-only", "仅训练 Adapter", "检验跨表示空间映射是否可学习", CYAN),
        ("U-Net-only", "仅微调 Lotus U-Net", "检验生成主干能否吸收 thermal + text 条件", GREEN),
        ("Adapter + U-Net", "联合训练", "检验更大自由度是否真的改善几何", YELLOW),
    ]
    for i, (name, tag, desc, color) in enumerate(cards):
        x = 0.72 + i * 3.12
        add_rect(s, x, 1.65, 2.82, 3.68, NAVY_2, True, GRID)
        add_rect(s, x, 1.65, 2.82, 0.09, color)
        add_text(s, name, x + 0.18, 1.98, 2.45, 0.42, 18, WHITE, True)
        add_text(s, tag, x + 0.18, 2.52, 2.45, 0.32, 12, color, True)
        add_text(s, desc, x + 0.18, 3.15, 2.42, 1.05, 13, MUTED)
        add_text(s, "同一 Val manifest\n同一 evaluator\n同一可视化", x + 0.18, 4.4, 2.3, 0.68, 11, WHITE, True)
    add_text(s, "决策标准：几何指标优先；diffusion loss 只作为训练诊断，不能替代深度质量。", 0.75, 5.85, 11.8, 0.5, 17, YELLOW, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 4 — no-caption route selection
    s = prs.slides.add_slide(blank); set_background(s); add_title(s, "无 Caption：没有单一赢家，但 Direct 明显落后", "Val · 5,810 samples", 4)
    columns = ["Metric", "Direct", "Adapter-only", "U-Net-only", "Adapter+U-Net"]
    rows = [
        ["AbsRel ↓", "0.4142", "0.3068", "0.2991", "0.2946"],
        ["SqRel ↓", "4.4449", "3.5674", "4.3983", "5.0678"],
        ["RMSE (m) ↓", "11.6435", "9.6402", "10.8475", "12.1988"],
        ["RMSElog ↓", "0.5818", "0.4128", "0.4026", "0.4120"],
        ["δ1 ↑", "0.3202", "0.4888", "0.5187", "0.5526"],
        ["δ2 ↑", "0.5864", "0.7695", "0.7905", "0.7947"],
        ["δ3 ↑", "0.7647", "0.8939", "0.9011", "0.8963"],
    ]
    highlights = {(0,4),(1,2),(2,2),(3,3),(4,4),(5,4),(6,3)}
    add_table(s, 0.65, 1.55, 12.0, 4.7, columns, rows, [2.15,2.0,2.2,2.2,2.55], 12, highlight_cells=highlights)
    add_text(s, "Adapter-only 最稳绝对误差；Joint 最强 AbsRel/δ1/δ2；U-Net-only 最均衡。", 0.7, 6.48, 12.0, 0.37, 15, WHITE, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 5 — correct captions across routes
    s = prs.slides.add_slide(blank); set_background(s); add_title(s, "加入 Correct Caption 后，路线响应并不一致", "Val · correct captions", 5)
    rows = [
        ["AbsRel ↓", "0.4135", "0.2834", "0.2970", "0.2918"],
        ["SqRel ↓", "4.4297", "3.4075", "4.3532", "4.7246"],
        ["RMSE (m) ↓", "11.6277", "9.5389", "10.7780", "11.8177"],
        ["RMSElog ↓", "0.5806", "0.3708", "0.3963", "0.4050"],
        ["δ1 ↑", "0.3204", "0.5127", "0.5232", "0.5428"],
        ["δ2 ↑", "0.5867", "0.7988", "0.7953", "0.7918"],
        ["δ3 ↑", "0.7650", "0.9196", "0.9042", "0.9006"],
    ]
    highlights = {(0,2),(1,2),(2,2),(3,2),(4,4),(5,2),(6,2)}
    add_table(s, 0.65, 1.55, 12.0, 4.55, columns, rows, [2.15,2.0,2.2,2.2,2.55], 12, highlight_cells=highlights)
    add_rich_text(s, [("关键点：", CYAN, True), ("“带 Caption 的模型更好”并不等于 Caption 有贡献；必须固定 checkpoint 做 Correct vs Empty 配对。", WHITE, False)], 0.75, 6.35, 11.8, 0.43, 14)
    add_footer(s)

    # 6 — paired val summary
    s = prs.slides.add_slide(blank); set_background(s); add_title(s, "同 checkpoint 配对：U-Net-only 给出最干净的 Caption 证据", "Val · causal comparison", 6)
    summary = [
        ("Direct", "≈ 0", "几乎不敏感", MUTED),
        ("Adapter-only", "Trade-off", "δ1 / RMSElog ↑；RMSE / SqRel ↓", YELLOW),
        ("U-Net-only", "7 / 7", "核心指标全部改善", GREEN),
        ("Adapter + U-Net", "0 / 7", "Correct 全面弱于 Empty", RED),
    ]
    for i, (name, score, note, color) in enumerate(summary):
        y = 1.5 + i * 1.18
        add_rect(s, 0.75, y, 11.8, 0.92, NAVY_2, True, GRID)
        add_text(s, name, 1.0, y + 0.18, 2.4, 0.35, 16, WHITE, True)
        add_text(s, score, 3.75, y + 0.12, 1.8, 0.42, 22, color, True, PP_ALIGN.CENTER)
        add_text(s, note, 6.0, y + 0.2, 5.95, 0.34, 14, MUTED)
    add_text(s, "因此后续 Caption 实验冻结 U-Net-only checkpoint，并进入独立 Test。", 0.75, 6.36, 11.8, 0.42, 17, GREEN, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 7 — U-Net val paired details
    s = prs.slides.add_slide(blank); set_background(s); add_title(s, "Val：Correct Caption 对 U-Net-only 的提升稳定且一致", "Paired evidence · 2,000 bootstrap", 7)
    val_rows = [
        ["AbsRel ↓", "0.2970", "0.3282", "85.6%", "+0.0313 [0.0305, 0.0321]"],
        ["SqRel ↓", "4.3532", "4.6104", "61.4%", "+0.2571 [0.2309, 0.2854]"],
        ["RMSE (m) ↓", "10.7780", "11.0377", "60.9%", "+0.2597 [0.2298, 0.2892]"],
        ["RMSElog ↓", "0.3963", "0.4469", "91.6%", "+0.0506 [0.0496, 0.0516]"],
        ["δ1 ↑", "0.5232", "0.4628", "87.8%", "+0.0604 [0.0591, 0.0617]"],
        ["δ2 ↑", "0.7953", "0.7478", "87.1%", "+0.0474 [0.0464, 0.0485]"],
        ["δ3 ↑", "0.9042", "0.8760", "88.5%", "+0.0281 [0.0275, 0.0288]"],
    ]
    add_table(s, 0.58, 1.5, 12.18, 4.95, ["Metric", "Correct", "Empty", "Correct win", "Mean improvement [95% CI]"], val_rows,
              [2.0,1.65,1.65,2.0,4.88], 11, highlight_cols={1})
    add_text(s, "Val 上 7 个核心指标的配对 95% CI 均不跨 0。", 0.75, 6.55, 11.8, 0.3, 15, GREEN, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 8 — frozen test
    s = prs.slides.add_slide(blank); set_background(s); add_title(s, "Frozen Test：Caption 贡献可以泛化，但存在平方误差取舍", "Independent test · 9,508 samples", 8)
    test_rows = [
        ["AbsRel ↓", "0.3492", "0.3616", "65.5%", "+0.0125"],
        ["SqRel ↓", "6.1464", "5.6157", "38.0%", "−0.5307"],
        ["RMSE (m) ↓", "13.0530", "12.7829", "42.3%", "−0.2701"],
        ["RMSElog ↓", "0.4679", "0.5013", "77.1%", "+0.0334"],
        ["δ1 ↑", "0.4627", "0.4186", "79.4%", "+0.0441"],
        ["δ2 ↑", "0.7255", "0.6967", "73.3%", "+0.0287"],
        ["δ3 ↑", "0.8517", "0.8345", "73.1%", "+0.0172"],
    ]
    add_table(s, 0.58, 1.5, 8.1, 4.95, ["Metric", "Correct", "Empty", "Win", "Δ"], test_rows,
              [1.75,1.55,1.55,1.35,1.9], 11, highlight_cols={1})
    add_metric_card(s, 9.02, 1.55, 3.45, 1.32, "RELATIVE / LOG METRICS", "稳定提升", "AbsRel · RMSElog · δ1/δ2/δ3", GREEN)
    add_metric_card(s, 9.02, 3.08, 3.45, 1.32, "SQUARED ERROR", "出现取舍", "SqRel 与 RMSE (m) 变差", RED)
    add_metric_card(s, 9.02, 4.61, 3.45, 1.32, "严谨表述", "误差尾部恶化", "不能仅凭聚合指标断言“远端变差”", YELLOW)
    add_footer(s, "Frozen Test · no checkpoint reselection after viewing results")

    # 9 — official route visual comparison
    s = prs.slides.add_slide(blank); set_background(s); add_title(s, "同一 Val 样本：四条路线的官方输出", "Official Iris/Lotus visualization", 9)
    seq = "2021-08-06-11-23-45"
    sample = "000000"
    route_images = [
        ("Direct", "direct_baseline_official_val_full"),
        ("Adapter-only", "adapter_only_official_val_full"),
        ("U-Net-only", "unet_only_official_val_full"),
        ("Adapter + U-Net", "adapter_unet_official_val_full"),
    ]
    for i, (label, directory) in enumerate(route_images):
        col, row = i % 2, i // 2
        x = 0.65 + col * 6.14
        y = 1.42 + row * 2.66
        add_text(s, label, x, y, 5.88, 0.28, 13, WHITE, True, PP_ALIGN.CENTER)
        # Source files are 640x256.  Use a matching wide box so the slide does
        # not introduce misleading black letterboxing around the official PNG.
        add_picture_contain(s, route_vis_path(directory, seq, sample), x, y + 0.34, 5.88, 2.35)
    add_text(s, "sample 000000 · untouched official vis/*.png · 预测图本身仍可能模糊", 0.62, 6.88, 12.0, 0.22, 10, MUTED, False, PP_ALIGN.CENTER)
    add_footer(s)

    # 10 — caption visual comparisons, successful + trade-off
    s = prs.slides.add_slide(blank); set_background(s); add_title(s, "官方可视化案例：成功样本与指标取舍样本", "U-Net-only · frozen Test", 10)
    correct_dir = "unet_only_correct_caption_official_test_full"
    empty_dir = "unet_only_caption_model_empty_official_test_full"
    seq = "2021-08-06-11-37-46"
    ids = [("003530", "成功例：AbsRel / δ1 / RMSE 同时改善", GREEN),
           ("009229", "取舍例：AbsRel / δ1 改善，但 RMSE 变差", YELLOW)]
    correct_csv = ROOT / "outputs" / "lotus_line_v1" / correct_dir / "per_sample_metrics.csv"
    empty_csv = ROOT / "outputs" / "lotus_line_v1" / empty_dir / "per_sample_metrics.csv"
    for row_i, (sid, desc, color) in enumerate(ids):
        y = 1.45 + row_i * 2.63
        cm = metrics_for_sample(correct_csv, sid); em = metrics_for_sample(empty_csv, sid)
        add_text(s, f"#{sid}", 0.52, y + 0.1, 0.75, 0.3, 12, color, True)
        add_text(s, desc, 0.52, y + 0.48, 2.1, 0.75, 11, WHITE, True)
        add_text(s, "EMPTY", 2.75, y - 0.02, 4.35, 0.26, 10, MUTED, True, PP_ALIGN.CENTER)
        add_picture_contain(s, route_vis_path(empty_dir, seq, sid), 2.75, y + 0.3, 4.35, 2.12)
        add_text(s, "CORRECT", 7.32, y - 0.02, 4.35, 0.26, 10, GREEN, True, PP_ALIGN.CENTER)
        add_picture_contain(s, route_vis_path(correct_dir, seq, sid), 7.32, y + 0.3, 4.35, 2.12)
        delta_abs = em["abs_relative_difference"] - cm["abs_relative_difference"]
        delta_d1 = cm["delta1_acc"] - em["delta1_acc"]
        delta_rmse = em["rmse_linear"] - cm["rmse_linear"]
        add_text(s, f"ΔAbsRel {delta_abs:+.3f}  ·  Δδ1 {delta_d1:+.3f}  ·  ΔRMSE {delta_rmse:+.2f} m",
                 11.88, y + 0.62, 1.0, 1.2, 10, color, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_footer(s, "Untouched official vis/*.png · same checkpoint / sample / evaluator")

    # 11 — conclusion
    s = prs.slides.add_slide(blank); set_background(s); add_title(s, "阶段性结论与下一步", "Takeaways", 11)
    add_metric_card(s, 0.75, 1.55, 3.65, 1.48, "01 · ROUTE", "选择 U-Net-only", "Caption 响应最一致、证据最干净", GREEN)
    add_metric_card(s, 4.82, 1.55, 3.65, 1.48, "02 · CONTRIBUTION", "Caption 有效", "Test 上相对 / 对数 / 阈值指标稳定提升", BLUE)
    add_metric_card(s, 8.89, 1.55, 3.65, 1.48, "03 · LIMITATION", "并非全面胜出", "SqRel / RMSE 暴露平方误差尾部问题", YELLOW)
    add_bullets(s, [
        "对老师：可以作为阶段性路线验证展示，结论有完整 Val 与独立 Test 支撑。",
        "对论文：后续冻结 manifest/config，按统一 MS2 protocol 重新评估最终 checkpoint。",
        "对诊断：若要定位 RMSE 恶化来源，需单独保存 raw prediction，做误差分位数与深度分箱；不要从聚合指标直接推断远端。",
        "对 Caption：下一阶段再设计鲁棒性与 hard-wrong 实验，当前不混入路线选择表。",
    ], 0.9, 3.55, 11.6, 2.45, 15, GREEN)
    add_text(s, "一句话：这条线已经从“能跑”推进到“可比较、可复现、可解释”。", 0.75, 6.38, 11.8, 0.44, 18, WHITE, True, PP_ALIGN.CENTER)
    add_footer(s)

    prs.save(OUT_PATH)
    print(OUT_PATH)
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    build()
