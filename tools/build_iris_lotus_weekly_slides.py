"""Build the 23-29 June 2026 Iris/Lotus weekly presentation.

The deck is deliberately self-contained and uses only verified local experiment
artifacts.  Run with the PowerShell ``pytorch_`` conda environment.
"""

from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "iris_lotus_weekly_2026-06-23_to_2026-06-29.pptx"

IMG_IRIS = ROOT / "img" / "iris.png"
IMG_BRIDGE = ROOT / "outputs" / "anythermal_lotus_direct_smoke" / "predicted_depth_visualization.png"
IMG_THERMAL = ROOT / "outputs" / "anythermal_lotus_direct_smoke" / "thermal_raw_visualization.png"
IMG_ADAPTER = ROOT / "outputs" / "adapter_v0_thermal_only_short_run_v2" / "plots" / "training_overview.png"
IMG_OVERFIT = ROOT / "outputs" / "adapter_v0_overfit32" / "loss_curve.png"
IMG_JOINT = ROOT / "outputs" / "adapter_v0_unet_joint_short_run" / "plots" / "fixed_val_visual_comparison_adapter_only_vs_joint.png"
IMG_PIXEL = ROOT / "outputs" / "adapter_v0_unet_joint_short_run" / "pixel_comparison_with_adapter_only" / "pixel_error_comparison_panel.png"
IMG_GROUND = ROOT / "outputs" / "iris_rgb_grounding_audit" / "visualizations" / "00_2021-08-06-11-23-45_000000.png"

SW, SH = Inches(13.333), Inches(7.5)

NAVY = RGBColor(16, 35, 55)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(100, 116, 139)
BG = RGBColor(247, 249, 252)
WHITE = RGBColor(255, 255, 255)
BORDER = RGBColor(218, 225, 234)
BLUE = RGBColor(47, 107, 255)
BLUE_PALE = RGBColor(232, 239, 255)
ORANGE = RGBColor(245, 158, 11)
ORANGE_PALE = RGBColor(255, 244, 222)
PURPLE = RGBColor(124, 58, 237)
PURPLE_PALE = RGBColor(242, 235, 255)
GREEN = RGBColor(22, 163, 74)
GREEN_PALE = RGBColor(229, 247, 235)
RED = RGBColor(220, 38, 38)
RED_PALE = RGBColor(254, 232, 232)
GRAY = RGBColor(107, 114, 128)
GRAY_PALE = RGBColor(238, 241, 245)

FONT = "Microsoft YaHei"
FONT_EN = "Aptos"


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def text_box(slide, text, x, y, w, h, size=18, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=FONT,
             margin=0.04, line_spacing=1.0):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def rich_text(slide, parts, x, y, w, h, size=18, align=PP_ALIGN.LEFT,
              valign=MSO_ANCHOR.MIDDLE):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for text, color, bold in parts:
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return shape


def box(slide, x, y, w, h, fill=WHITE, line=BORDER, radius=True,
        line_width=1.0):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_width)
    return shp


def labeled_box(slide, title, subtitle, x, y, w, h, fill, accent,
                title_size=16, subtitle_size=10):
    box(slide, x, y, w, h, fill=fill, line=accent)
    text_box(slide, title, x + 0.12, y + 0.08, w - 0.24, 0.32,
             size=title_size, color=accent, bold=True, valign=MSO_ANCHOR.MIDDLE)
    if subtitle:
        text_box(slide, subtitle, x + 0.12, y + 0.41, w - 0.24, h - 0.48,
                 size=subtitle_size, color=MUTED, valign=MSO_ANCHOR.TOP)


def chip(slide, text, x, y, w, color, pale):
    box(slide, x, y, w, 0.4, fill=pale, line=color)
    text_box(slide, text, x + 0.04, y + 0.03, w - 0.08, 0.32,
             size=11, color=color, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)


def arrow(slide, x, y, w=0.38, h=0.3, color=MUTED, direction="right"):
    kind = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.DOWN_ARROW
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def title(slide, heading, kicker=None, n=None):
    if kicker:
        text_box(slide, kicker.upper(), 0.55, 0.25, 6.5, 0.25,
                 size=9, color=BLUE, bold=True, font=FONT_EN)
    text_box(slide, heading, 0.55, 0.48, 12.1, 0.55,
             size=26, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.04), Inches(12.2), Inches(0.025))
    line.fill.solid(); line.fill.fore_color.rgb = BORDER; line.line.fill.background()
    if n is not None:
        text_box(slide, f"{n:02d}", 12.43, 0.28, 0.35, 0.25,
                 size=9, color=MUTED, bold=True, align=PP_ALIGN.RIGHT, font=FONT_EN)


def footer(slide, source):
    text_box(slide, source, 0.58, 7.16, 11.8, 0.17, size=7.5, color=MUTED,
             font=FONT_EN, valign=MSO_ANCHOR.MIDDLE)


def add_image_contain(slide, path: Path, x, y, w, h, border=True, bg=WHITE):
    if not path.exists():
        raise FileNotFoundError(path)
    if border:
        box(slide, x, y, w, h, fill=bg, line=BORDER, radius=False)
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))


def set_cell(cell, text, size=11, color=TEXT, bold=False, align=PP_ALIGN.CENTER,
             fill=WHITE):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.margin_left = cell.margin_right = Inches(0.06)
    cell.margin_top = cell.margin_bottom = Inches(0.03)
    tf = cell.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color


def simple_table(slide, rows, x, y, w, h, widths=None, header_fill=NAVY,
                 header_size=10.5, body_size=10.5, best=None):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    if widths:
        for col, frac in zip(table.columns, widths):
            col.width = Inches(w * frac)
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            if i == 0:
                set_cell(table.cell(i, j), str(value), size=header_size, color=WHITE,
                         bold=True, fill=header_fill)
            else:
                is_best = best and (i, j) in best
                fill = BLUE_PALE if is_best else (WHITE if i % 2 else RGBColor(248, 250, 252))
                set_cell(table.cell(i, j), str(value), size=body_size,
                         color=BLUE if is_best else TEXT, bold=is_best, fill=fill)
    return table


def takeaway(slide, text, y=6.58, color=NAVY, pale=WHITE):
    box(slide, 0.68, y, 11.98, 0.48, fill=pale, line=color)
    rich_text(slide, [("结论｜", color, True), (text, TEXT, True)],
              0.82, y + 0.04, 11.65, 0.36, size=13)


def add_bullets(slide, bullets, x, y, w, h, size=14, color=TEXT,
                bullet_color=BLUE, gap=0.05):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap * 72)
        p.line_spacing = 1.05
        p.level = 0
        p.text = "•  " + bullet
        for r in p.runs:
            r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = color
    return shape


def build_deck():
    required = [IMG_IRIS, IMG_BRIDGE, IMG_THERMAL, IMG_ADAPTER, IMG_OVERFIT,
                IMG_JOINT, IMG_PIXEL, IMG_GROUND]
    missing = [str(p) for p in required if not p.exists()]
    if missing:
        raise FileNotFoundError("Missing required slide assets:\n" + "\n".join(missing))

    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    blank = prs.slide_layouts[6]

    # Slide 1
    s = prs.slides.add_slide(blank); set_bg(s, WHITE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    text_box(s, "IRIS / LOTUS ROUTE", 0.72, 0.7, 4.8, 0.3, size=10,
             color=BLUE, bold=True, font=FONT_EN)
    text_box(s, "Iris/Lotus Reproduction\nand Grounding Audit", 0.72, 1.15, 5.85, 1.65,
             size=30, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE, font=FONT_EN)
    text_box(s, "语言条件扩散深度估计的复现、\n热成像适配与能力验证", 0.75, 3.02, 5.4, 0.88,
             size=20, color=TEXT, bold=True)
    chip(s, "Weekly Progress", 0.77, 4.25, 1.65, BLUE, BLUE_PALE)
    text_box(s, "23–29 June 2026", 2.56, 4.26, 2.3, 0.36, size=12,
             color=MUTED, font=FONT_EN, valign=MSO_ANCHOR.MIDDLE)
    rich_text(s, [("复现接口", BLUE, True), ("  →  ", MUTED, False),
                  ("学习适配", ORANGE, True), ("  →  ", MUTED, False),
                  ("能力审计", GREEN, True)],
              0.73, 5.2, 5.7, 0.55, size=14)
    box(s, 6.75, 0.55, 5.95, 6.15, fill=BG, line=BORDER)
    add_image_contain(s, IMG_IRIS, 6.95, 0.78, 5.55, 3.05, border=False, bg=BG)
    add_image_contain(s, IMG_GROUND, 6.95, 4.05, 5.55, 2.30, border=False, bg=BG)
    text_box(s, "RGB + Caption → diffusion denoising → depth", 7.05, 6.35, 5.3, 0.2,
             size=8, color=MUTED, align=PP_ALIGN.CENTER, font=FONT_EN)

    # Slide 2
    s = prs.slides.add_slide(blank); set_bg(s); title(s, "Iris：用语言缩小单目深度的解空间", "Method", 2)
    labeled_box(s, "单目歧义", "同一二维观测\n可能对应多个三维解释", 0.62, 1.35, 2.15, 1.02, GRAY_PALE, GRAY)
    labeled_box(s, "困难区域", "小物体 · 遮挡\n弱纹理 · 尺度不确定", 0.62, 2.55, 2.15, 1.02, RED_PALE, RED)
    labeled_box(s, "语言约束", "物体身份 · 大小\n空间关系 · 场景结构", 0.62, 3.75, 2.15, 1.02, GREEN_PALE, GREEN)
    text_box(s, "训练：Caption 直接进入 denoising", 3.18, 1.28, 4.6, 0.35,
             size=14, color=NAVY, bold=True)
    xs = [3.18, 4.78, 6.38, 8.13, 10.0]
    labels = [
        ("RGB image", "x", BLUE, BLUE_PALE),
        ("Frozen VAE", "image latent", GRAY, GRAY_PALE),
        ("Noisy depth", "zₜ", PURPLE, PURPLE_PALE),
        ("Diffusion U-Net", "+ text tokens", PURPLE, PURPLE_PALE),
        ("Depth loss", "noise / x₀", ORANGE, ORANGE_PALE),
    ]
    for idx, (lab, sub, c, pale) in enumerate(labels):
        labeled_box(s, lab, sub, xs[idx], 1.78, 1.34 if idx != 3 else 1.55, 0.9, pale, c, 13, 9)
        if idx < len(labels)-1:
            arrow(s, xs[idx] + (1.34 if idx != 3 else 1.55) + 0.06, 2.07, 0.34, 0.25, MUTED)
    labeled_box(s, "Caption c", "Frozen CLIP\nper-token states", 6.95, 3.15, 1.75, 1.0, GREEN_PALE, GREEN)
    arrow(s, 7.66, 2.75, 0.28, 0.34, GREEN, direction="down")
    text_box(s, "推理：Gaussian depth latent 经过迭代去噪", 3.18, 4.5, 5.2, 0.34,
             size=14, color=NAVY, bold=True)
    infer = [("Gaussian zₜ", PURPLE, PURPLE_PALE), ("RGB + Caption", GREEN, GREEN_PALE),
             ("Iterative denoising", PURPLE, PURPLE_PALE), ("VAE Decoder", GRAY, GRAY_PALE),
             ("Predicted depth", BLUE, BLUE_PALE)]
    ix = 3.2
    for i, (lab, c, pale) in enumerate(infer):
        labeled_box(s, lab, "", ix, 5.0, 1.58, 0.62, pale, c, 12, 9)
        if i < len(infer)-1:
            arrow(s, ix + 1.65, 5.17, 0.34, 0.24, MUTED)
        ix += 2.02
    takeaway(s, "文本不是后处理，而是 U-Net 每一步去噪时使用的条件。", 6.36, GREEN, GREEN_PALE)
    footer(s, "Source: README.md; lotus/train_iris_g.py; lotus/pipeline.py")

    # Slide 3
    s = prs.slides.add_slide(blank); set_bg(s); title(s, "Iris 并不是单一模型，而是一套文本接入策略", "Repository map", 3)
    rows = [
        ["模型 / 路线", "主要任务与特点"],
        ["Lotus-D", "Depth-only dense prediction"],
        ["Lotus-G", "RGB condition + noisy depth latent"],
        ["Marigold", "Iterative latent diffusion depth"],
        ["E2E-FT", "端到端扩散深度微调"],
        ["Iris", "为上述模型增加语言条件的策略"],
    ]
    simple_table(s, rows, 0.65, 1.38, 5.2, 3.85, widths=[0.32, 0.68], body_size=11)
    text_box(s, "为什么聚焦 Lotus-G？", 0.77, 5.52, 2.6, 0.35, size=14, color=NAVY, bold=True)
    add_bullets(s, ["VAE latent 接口明确", "图像 / depth / text 三路可独立审计", "适合验证 AnyThermal 条件替换"],
                0.75, 5.9, 4.9, 0.95, size=11.5)
    text_box(s, "本地代码关系", 6.35, 1.37, 2.4, 0.35, size=14, color=NAVY, bold=True)
    labeled_box(s, "pipeline / inference", "LotusGPipeline", 7.28, 1.85, 3.3, 0.85, BLUE_PALE, BLUE)
    arrow(s, 8.72, 2.78, 0.32, 0.38, MUTED, direction="down")
    labeled_box(s, "Lotus U-Net + VAE + Scheduler", "Stable Diffusion latent space", 7.28, 3.22, 3.3, 0.85, PURPLE_PALE, PURPLE)
    arrow(s, 8.72, 4.15, 0.32, 0.38, MUTED, direction="down")
    labeled_box(s, "train_iris_d.py / train_iris_g.py", "CLIP encoder_hidden_states", 7.28, 4.58, 3.3, 0.85, GREEN_PALE, GREEN)
    box(s, 6.35, 5.75, 6.15, 0.72, fill=WHITE, line=BORDER)
    rich_text(s, [("边界：", NAVY, True), (" Iris ≠ Lotus ≠ Stable Diffusion；", TEXT, False),
                  ("Iris 是条件策略。", BLUE, True)], 6.55, 5.9, 5.75, 0.4, size=13)
    takeaway(s, "先把模型边界和条件流向说清楚，后续实验才知道在验证哪一层。", 6.55)
    footer(s, "Source: README.md; train_iris_d.py; train_iris_g.py")

    # Slide 4
    s = prs.slides.add_slide(blank); set_bg(s); title(s, "Zero-parameter Bridge：Shape 对齐不等于语义对齐", "Interface failure", 4)
    flow = [
        ("Thermal", "256×640", BLUE, BLUE_PALE, 1.35),
        ("Frozen AnyThermal", "[B,768,18,45]", BLUE, BLUE_PALE, 1.85),
        ("Grouped average", "768 → 4", ORANGE, ORANGE_PALE, 1.75),
        ("Resize", "[B,4,32,80]", GRAY, GRAY_PALE, 1.35),
        ("Lotus U-Net", "forward ✓", PURPLE, PURPLE_PALE, 1.55),
    ]
    x = 0.65
    for i, (lab, sub, c, pale, ww) in enumerate(flow):
        labeled_box(s, lab, sub, x, 1.35, ww, 0.78, pale, c, 12, 9)
        if i < len(flow)-1:
            arrow(s, x + ww + 0.06, 1.61, 0.32, 0.23, MUTED)
        x += ww + 0.48
    text_box(s, "实测表示分布", 0.68, 2.52, 2.1, 0.32, size=14, color=NAVY, bold=True)
    rows = [
        ["表示", "Shape", "Mean", "Std"],
        ["Zero-param bridge", "4×32×80", "−0.0261", "0.0979"],
        ["RGB VAE latent", "4×30×96", "−0.0166", "0.8493"],
    ]
    simple_table(s, rows, 0.65, 2.92, 6.15, 1.55, widths=[0.36, 0.24, 0.2, 0.2], body_size=10.5)
    box(s, 0.65, 4.72, 6.15, 1.32, fill=WHITE, line=BORDER)
    rich_text(s, [("8.7×", RED, True), ("  std 差距", NAVY, True)], 0.9, 4.88, 1.85, 0.45, size=24)
    text_box(s, "AnyThermal 是 DINOv2 patch representation；\nLotus 期待 Stable Diffusion VAE latent。", 2.75, 4.82, 3.72, 0.82,
             size=13, color=TEXT, bold=True, valign=MSO_ANCHOR.MIDDLE)
    text_box(s, "直接接入后的典型输出", 7.2, 2.52, 2.6, 0.32, size=14, color=NAVY, bold=True)
    add_image_contain(s, IMG_BRIDGE, 7.18, 2.92, 5.45, 3.15)
    chip(s, "中心黑 / 周围白 / 环状伪影", 8.45, 5.62, 3.0, RED, RED_PALE)
    takeaway(s, "forward 成功只证明 shape 可用；环状伪影暴露的是 latent space mismatch。", 6.48, RED, RED_PALE)
    footer(s, "Source: outputs/anythermal_lotus_direct_smoke/run_info.json; docs/iris_lotus_image_condition_interface.md")

    # Slide 5
    s = prs.slides.add_slide(blank); set_bg(s); title(s, "用可学习 Adapter 建立 AnyThermal → Lotus 接口", "Learned mapping", 5)
    levels = [("L8", 1.55), ("L9", 2.32), ("L10", 3.09), ("L11", 3.86)]
    for lab, yy in levels:
        labeled_box(s, lab, "768×18×45", 0.65, yy, 1.25, 0.58, BLUE_PALE, BLUE, 11, 8)
        arrow(s, 1.98, yy + 0.17, 0.32, 0.22, MUTED)
    labeled_box(s, "Per-level projection", "1×1 Conv 768→128 + GELU", 2.4, 2.08, 2.05, 1.18, ORANGE_PALE, ORANGE, 13, 9)
    arrow(s, 4.58, 2.52, 0.38, 0.28, MUTED)
    labeled_box(s, "Resize + Fusion", "concat 512 → 128", 5.05, 2.08, 1.75, 1.18, ORANGE_PALE, ORANGE, 13, 9)
    arrow(s, 6.92, 2.52, 0.38, 0.28, MUTED)
    labeled_box(s, "4-ch condition", "GroupNorm + scale/bias\n[B,4,32,80]", 7.38, 2.02, 1.75, 1.3, ORANGE_PALE, ORANGE, 13, 9)
    arrow(s, 9.25, 2.52, 0.38, 0.28, MUTED)
    labeled_box(s, "Lotus U-Net", "+ noisy depth latent\n[B,8,32,80]", 9.72, 2.02, 1.7, 1.3, PURPLE_PALE, PURPLE, 13, 9)
    chip(s, "AnyThermal frozen", 0.65, 4.72, 1.75, GRAY, GRAY_PALE)
    chip(s, "Adapter trainable", 2.55, 4.72, 1.75, ORANGE, ORANGE_PALE)
    chip(s, "Lotus U-Net frozen", 4.45, 4.72, 1.9, GRAY, GRAY_PALE)
    text_box(s, "32-sample overfit", 7.02, 4.55, 2.0, 0.3, size=12, color=NAVY, bold=True, font=FONT_EN)
    add_image_contain(s, IMG_OVERFIT, 6.98, 4.9, 3.35, 1.38)
    box(s, 10.58, 4.62, 2.0, 1.5, fill=WHITE, line=ORANGE)
    text_box(s, "984,212", 10.78, 4.78, 1.6, 0.45, size=23, color=ORANGE, bold=True,
             align=PP_ALIGN.CENTER, font=FONT_EN)
    text_box(s, "trainable parameters", 10.78, 5.22, 1.6, 0.25, size=9, color=MUTED,
             align=PP_ALIGN.CENTER, font=FONT_EN)
    rich_text(s, [("0.2931", RED, True), ("  →  ", MUTED, False), ("0.0747", GREEN, True)],
              10.73, 5.55, 1.7, 0.35, size=14, align=PP_ALIGN.CENTER)
    takeaway(s, "Adapter 能学习跨表示映射；32 张过拟合是可学习性验证，不是泛化结果。", 6.52, ORANGE, ORANGE_PALE)
    footer(s, "Source: models/anythermal_lotus_adapter.py; outputs/adapter_v0_overfit32/summary.json")

    # Slide 6
    s = prs.slides.add_slide(blank); set_bg(s); title(s, "Diffusion Loss 更低，不代表深度几何更好", "Optimization audit", 6)
    rows = [
        ["训练方式", "Val diffusion loss ↓", "Corr ↑", "Aligned RMSE ↓"],
        ["Adapter-only", "0.26965", "0.75975", "0.09713"],
        ["Adapter + U-Net", "0.23162", "0.69925", "0.10627"],
    ]
    simple_table(s, rows, 0.65, 1.35, 8.0, 1.42, widths=[0.28, 0.27, 0.2, 0.25],
                 body_size=10.5, best={(2, 1), (1, 2), (1, 3)})
    box(s, 8.92, 1.35, 3.72, 1.42, fill=WHITE, line=BORDER)
    text_box(s, "固定 8 张样本", 9.12, 1.52, 1.8, 0.28, size=12, color=NAVY, bold=True)
    rich_text(s, [("0.11828", BLUE, True), ("  vs  ", MUTED, False), ("0.13038", RED, True)],
              9.1, 1.88, 3.3, 0.42, size=20, align=PP_ALIGN.CENTER)
    text_box(s, "Aligned RMSE ↓   Adapter-only / Joint", 9.1, 2.34, 3.3, 0.2,
             size=8.5, color=MUTED, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_image_contain(s, IMG_JOINT, 0.65, 3.03, 11.98, 3.12)
    chip(s, "Adapter-only", 0.88, 5.73, 1.45, BLUE, BLUE_PALE)
    chip(s, "Joint U-Net", 2.48, 5.73, 1.35, PURPLE, PURPLE_PALE)
    takeaway(s, "Joint 继续优化扩散目标，却没有转化为更好的像素级几何。", 6.47, PURPLE, PURPLE_PALE)
    footer(s, "Source: outputs/adapter_v0_unet_joint_short_run/summary.json and comparison_with_adapter_only.json")

    # Slide 7
    s = prs.slides.add_slide(blank); set_bg(s); title(s, "当前 RGB 模型是否真的利用 Caption？", "Grounding protocol", 7)
    box(s, 0.65, 1.32, 12.0, 1.15, fill=WHITE, line=BORDER)
    text_box(s, "受控变量", 0.85, 1.5, 1.0, 0.3, size=13, color=NAVY, bold=True)
    controls = ["16 RGB samples", "same checkpoint", "seed 20260629", "same scheduler", "same preprocessing"]
    cx = 2.0
    for i, c in enumerate(controls):
        chip(s, c, cx, 1.68, 1.76 if i != 2 else 1.65, BLUE, BLUE_PALE)
        cx += 1.94
    text_box(s, "唯一改变：Caption", 0.72, 2.82, 2.2, 0.35, size=14, color=NAVY, bold=True)
    caps = [
        ("correct", "与图像一致", GREEN, GREEN_PALE),
        ("empty", "空文本", GRAY, GRAY_PALE),
        ("generic", "A driving scene.", BLUE, BLUE_PALE),
        ("hard-wrong", "交换左右 / 远近", RED, RED_PALE),
    ]
    x = 0.72
    for lab, sub, c, pale in caps:
        labeled_box(s, lab, sub, x, 3.25, 2.62, 0.9, pale, c, 15, 9)
        x += 3.02
    text_box(s, "工程排除项", 0.72, 4.58, 2.0, 0.35, size=14, color=NAVY, bold=True)
    checks = [
        ("Checkpoint", "missing / unexpected / mismatch = 0"),
        ("CLIP states", "逐 token hidden states"),
        ("Cross-attention", "16 个 U-Net attn2 层实际调用"),
        ("Determinism", "same prompt + seed → max diff 0"),
    ]
    x = 0.72
    for lab, sub in checks:
        labeled_box(s, "✓  " + lab, sub, x, 5.05, 2.83, 0.85, GREEN_PALE, GREEN, 12, 8.5)
        x += 3.02
    takeaway(s, "先排除加载、接线和随机性，再判断学习到的 grounding 是否可靠。", 6.48, GREEN, GREEN_PALE)
    footer(s, "Source: outputs/iris_rgb_grounding_audit/checkpoint_load_report.json; audit_report.md")

    # Slide 8
    s = prs.slides.add_slide(blank); set_bg(s); title(s, "Text Interface 已接通，但 Correct Caption 没有稳定优势", "Grounding result", 8)
    add_image_contain(s, IMG_GROUND, 0.65, 1.28, 12.0, 3.02)
    rows = [
        ["对比", "Correct 胜率"],
        ["vs hard-wrong · Norm-MSE", "50.0%"],
        ["vs hard-wrong · AbsRel", "37.5%"],
        ["vs hard-wrong · RMSE", "25.0%"],
        ["vs hard-wrong · ≥3/4 metrics", "25.0%"],
        ["vs empty · Norm-MSE", "12.5%"],
    ]
    simple_table(s, rows, 0.65, 4.53, 6.25, 1.86, widths=[0.72, 0.28], header_fill=NAVY,
                 header_size=9.5, body_size=9.5)
    box(s, 7.18, 4.53, 5.47, 1.86, fill=WHITE, line=RED)
    rich_text(s, [("0.52%", RED, True), ("  correct ↔ hard-wrong\n", TEXT, True)],
              7.48, 4.72, 2.15, 0.72, size=21, valign=MSO_ANCHOR.TOP)
    text_box(s, "相对差异 RMSE / prediction std", 7.5, 5.33, 2.2, 0.35,
             size=9, color=MUTED, font=FONT_EN)
    rich_text(s, [("0.992", RED, True), ("  attention entropy", TEXT, True)],
              9.82, 4.72, 2.45, 0.5, size=20)
    text_box(s, "高度分散；未证明物体级聚焦", 9.84, 5.25, 2.35, 0.45,
             size=10, color=MUTED)
    text_box(s, "当前 checkpoint 会响应文本，但没有稳定的 token → region → geometry 对齐证据。",
             7.5, 5.78, 4.78, 0.42, size=12, color=RED, bold=True, align=PP_ALIGN.CENTER)
    takeaway(s, "Interface 正常；失败定位在已学习的语义 / 空间 grounding，而不是文本未进入 U-Net。", 6.54, RED, RED_PALE)
    footer(s, "Source: outputs/iris_rgb_grounding_audit/metrics_summary.json and visualizations/00_*.png")

    # Slide 9
    s = prs.slides.add_slide(blank); set_bg(s); title(s, "Iris/Lotus 路线的阶段性结论", "Takeaways", 9)
    columns = [
        ("已经确认", GREEN, GREEN_PALE, [
            "图像、深度与文本接口已梳理",
            "Learned Adapter 可以学习映射",
            "RGB text tokens 进入 cross-attention",
        ]),
        ("已经排除", ORANGE, ORANGE_PALE, [
            "Shape 一致 ≠ latent 兼容",
            "diffusion loss 更低 ≠ geometry 更好",
            "调用 attention ≠ 物体级 grounding",
        ]),
        ("当前限制", RED, RED_PALE, [
            "correct 未稳定优于 hard-wrong",
            "本地 checkpoint 不宜作为 teacher",
            "结论不外推为对论文整体的否定",
        ]),
    ]
    x = 0.65
    for heading, c, pale, bullets in columns:
        box(s, x, 1.42, 3.82, 3.65, fill=WHITE, line=c, line_width=1.5)
        box(s, x, 1.42, 3.82, 0.62, fill=pale, line=c)
        text_box(s, heading, x + 0.18, 1.56, 3.46, 0.3, size=16, color=c, bold=True,
                 align=PP_ALIGN.CENTER)
        add_bullets(s, bullets, x + 0.24, 2.28, 3.34, 2.2, size=13, bullet_color=c)
        x += 4.08
    box(s, 0.65, 5.42, 12.0, 1.18, fill=NAVY, line=NAVY)
    text_box(s, "从“尝试连接模型”推进到：", 0.95, 5.64, 3.1, 0.32,
             size=13, color=RGBColor(190, 210, 232), bold=True)
    rich_text(s, [("接口问题", BLUE, True), ("  /  ", WHITE, False),
                  ("latent mismatch", ORANGE, True), ("  /  ", WHITE, False),
                  ("优化目标", PURPLE, True), ("  /  ", WHITE, False),
                  ("文本 grounding", GREEN, True)],
              3.82, 5.58, 8.25, 0.48, size=18, align=PP_ALIGN.CENTER)
    text_box(s, "当前审计针对本地可复现 checkpoint，而非对作者完整 text-trained model 的最终否定。",
             1.05, 6.17, 11.0, 0.26, size=10, color=WHITE, align=PP_ALIGN.CENTER)
    footer(s, "Weekly Progress | Iris/Lotus route | 23–29 June 2026")

    assert len(prs.slides) == 9
    # Geometry check: permit tiny rounding noise only.
    for idx, slide in enumerate(prs.slides, 1):
        for shp in slide.shapes:
            if shp.left < 0 or shp.top < 0 or shp.left + shp.width > SW + 5 or shp.top + shp.height > SH + 5:
                raise ValueError(f"Slide {idx} shape out of bounds: {shp.name}")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    with ZipFile(OUT) as zf:
        if zf.testzip() is not None:
            raise ValueError("Corrupt PPTX ZIP structure")
        slide_xml = [n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        if len(slide_xml) != 9:
            raise ValueError(f"Expected 9 slide XML files, found {len(slide_xml)}")
    print(f"created={OUT}")
    print("slides=9")
    print(f"bytes={OUT.stat().st_size}")


if __name__ == "__main__":
    build_deck()
